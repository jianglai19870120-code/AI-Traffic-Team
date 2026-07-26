#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from __future__ import annotations

import argparse
from collections import Counter
import json
import pathlib
import re
import sys
from typing import Any

from PIL import Image
from pptx import Presentation

from common import (
    SKILL_ROOT,
    ensure_dir,
    load_character,
    load_config,
    resolve_default_character,
    slugify,
    write_json,
    write_text,
)


ODD_BG = "#fcb537"
EVEN_BG = "#122142"
ODD_LINE = "black"
EVEN_LINE = "white"
SKIN_TONE_BASE = "#e8a668"
COVER_BG = ODD_BG
DEFAULT_COVER_SUPPORT_LABELS = ["After Work", "Monetize", "Results"]
DEFAULT_COVER_ENGLISH_LABELS = ["After Work", "Monetize", "Results", "Income", "Leverage", "Value"]
COVER_ACCENT_BLOCKS_MIN = 1
COVER_ACCENT_BLOCKS_MAX = 3
FACE_ANCHOR_FILENAME = "ip-face-anchor-top-left.png"
ACTION_REFERENCE_BY_INTENT = [
    (("讲解", "推进", "说明", "流程", "操作", "演示"), ("ip-ref-02-speaking-mic.png", "ip-ref-04-questioning.png")),
    (("思考", "判断", "分析", "拆解"), ("ip-ref-03-thinking.png",)),
    (("强调", "阻挡", "误区", "提醒"), ("ip-ref-05-stop-pose.png",)),
    (("指向", "号召", "结论", "总结", "顿悟"), ("ip-ref-06-pointing.png", "ip-ref-09-idea.png")),
    (("震惊", "踩坑", "强提醒", "惊讶"), ("ip-ref-07-shocked.png",)),
    (("困惑", "卡点", "不理解", "质疑"), ("ip-ref-08-confused.png",)),
]
DEFAULT_LAYOUT_CONSTRAINTS = [
    "标题区优先",
    "说明文案小字号",
    "配图和配字分离",
    "文字原文锚定，只允许少量源自原文的结构词",
    "左下角 13cm × 13cm 纯底色禁绘区域",
    "左下角不要出现任何文字",
    "左下角不要出现人物、图标、线条、箭头、坑位边缘和装饰符号",
]

STEP_LABEL_RE = re.compile(
    r"^\s*(?P<label>(?:第[一二三四五六七八九十百千万0-9]+[点步])|(?:[一二三四五六七八九十]+、)|(?:\d+[\.、]))"
)
EXPLICIT_STEP_TITLE_RE = re.compile(
    r"(?:第[一二三四五六七八九十百千万0-9]+(?:层|步|点|条)|认知[一二三四五六七八九十百千万0-9]+)"
)
STEP_TITLE_ENTRY_RE = re.compile(
    r"(?P<entry>(?:第[一二三四五六七八九十百千万0-9]+(?:层|步|点|条)|认知[一二三四五六七八九十百千万0-9]+)(?:[：:，,、 ]*[^。！？；\n]+)?)"
)
GUIDE_TITLE_PREFIX_RE = re.compile(
    r"^\s*(?P<prefix>(?:第[一二三四五六七八九十百千万0-9]+(?:点|步|条|讲|课|章|节|部分|阶段))|(?:认知[一二三四五六七八九十百千万0-9]+)|(?:误区[一二三四五六七八九十百千万0-9]+)|(?:步骤[一二三四五六七八九十百千万0-9]+)|(?:方法[一二三四五六七八九十百千万0-9]+)|(?:模块[一二三四五六七八九十百千万0-9]+))"
)
SUMMARY_SENTENCE_RE = re.compile(r"^(所以|总之|最后|说到底|归根结底|真正厉害的人|记住)")
INTRO_BREAK_RE = re.compile(r"(今天我就|先别|接下来|你最该|真正该补的是)")
PAGE_BREAK_RE = re.compile(r"(所以|但是|因为|真正该补的是|正确顺序是|也就是说|我见过|尤其在)")
PAGE_REF_RE = re.compile(r"^\s*(\d+)(?:\s*-\s*(\d+))?\s*$")
DECK_TITLE_RE = re.compile(r"^\s*选题[：:]\s*(.+?)\s*$")
WRAPPER_HEADING_RE = re.compile(r"^\s*#+\s*.+生成结果\s*$")
WRAPPER_META_RE = re.compile(r"^\s*(调用结构编号|调用结构文件名|生成时间)[：:].*$")
INVALID_WINDOWS_NAME_RE = re.compile(r'[<>:"/\\|?*]+')
RESERVED_ZONE = {
    "position": "bottom-left",
    "size_cm": "13 x 13",
    "fill_mode": "pure-background-no-draw",
}
RESERVED_ZONE_RULES = [
    "no_text",
    "no_icons",
    "no_character",
    "no_lines",
    "no_shapes",
    "no_scene_overlap",
    "background_only",
    "no_frame",
    "no_border",
    "no_placeholder_box",
]
BOLD_TEXT_RE = re.compile(r"\*\*(.+?)\*\*")

COVER_OUTPUT_SPECS = [
    {
        "cover_type": "cover-3x4",
        "aspect_ratio": "3:4",
        "filename": "cover-3x4.png",
        "title_zone": "top-centered",
        "composition_guidance": "人物居中或偏下居中，标题在上方主标题区，标签围绕人物点缀",
    },
    {
        "cover_type": "cover-4x3",
        "aspect_ratio": "4:3",
        "filename": "cover-4x3.png",
        "title_zone": "left-or-top-horizontal",
        "composition_guidance": "人物偏一侧，标题与标签形成横向信息区，整体更像横版主视觉卡面",
    },
]


def unique_preserve_order(values: list[str]) -> list[str]:
    result: list[str] = []
    for value in values:
        normalized = value.strip()
        if normalized and normalized not in result:
            result.append(normalized)
    return result


def strip_markdown_bold(text: str) -> str:
    return BOLD_TEXT_RE.sub(lambda match: match.group(1).strip(), text).strip()


def extract_bold_title(line: str) -> str:
    matches = [match.strip() for match in BOLD_TEXT_RE.findall(line) if match.strip()]
    if not matches:
        return ""
    fully_bold = re.fullmatch(r"\*\*\s*(.+?)\s*\*\*", line.strip())
    if fully_bold:
        return fully_bold.group(1).strip()
    return " ".join(matches).strip()


def clean_input_text(text: str) -> str:
    lines: list[str] = []
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith("|") and line.endswith("|"):
            parts = [part.strip() for part in line.strip("|").split("|")]
            if not parts:
                continue
            if all(re.fullmatch(r"[-:\s]+", part or "") for part in parts):
                continue
            line = " ".join(part for part in parts if part)
        fragments = [fragment.strip() for fragment in re.split(r"<br\s*/?>", line, flags=re.IGNORECASE)]
        for fragment in fragments:
            if not fragment:
                continue
            if WRAPPER_HEADING_RE.match(fragment):
                continue
            if WRAPPER_META_RE.match(fragment):
                continue
            if re.fullmatch(r"-{3,}", fragment):
                lines.append("---")
                continue
            lines.append(fragment)
    return "\n".join(lines).strip()


def extract_deck_title(text: str) -> tuple[str | None, str]:
    normalized = clean_input_text(text)
    if not normalized:
        return None, ""
    lines = normalized.splitlines()
    if not lines:
        return None, normalized
    title_index: int | None = None
    deck_title: str | None = None
    for index, raw_line in enumerate(lines):
        line = raw_line.strip()
        if not line:
            continue
        match = DECK_TITLE_RE.match(line)
        if match:
            title_index = index
            deck_title = strip_markdown_bold(match.group(1).strip())
            break
        if deck_title is None:
            title_index = index
            deck_title = strip_markdown_bold(line)
            break
    if title_index is None:
        return None, normalized
    remaining_lines = [line.strip() for line in lines[title_index + 1 :] if line.strip()]
    while remaining_lines and DECK_TITLE_RE.match(remaining_lines[0]):
        remaining_lines.pop(0)
    remaining = "\n".join(remaining_lines).strip()
    return deck_title or None, remaining


def sanitize_windows_name(name: str) -> str:
    cleaned = INVALID_WINDOWS_NAME_RE.sub("", name).strip().rstrip(".")
    cleaned = re.sub(r"\s+", " ", cleaned)
    if cleaned:
        return cleaned
    return slugify(name)


def resolve_archive_root(cfg: dict[str, Any]) -> pathlib.Path:
    configured_root = str(cfg.get("archive_root") or "").strip()
    if not configured_root:
        raise ValueError(
            "archive_root 未显式配置，当前已禁止回退到 skill 相邻目录。"
            "请先在 config/skill-config.json 中设置正式归档根目录，再继续正式出图。"
        )
    return pathlib.Path(configured_root).expanduser().resolve()


def resolve_archive_dir(deck_title: str | None, input_path: pathlib.Path, archive_root: pathlib.Path) -> pathlib.Path:
    if not deck_title or not deck_title.strip():
        raise ValueError(
            f"输入文案缺少选题标题，正式归档子文件夹名只允许使用文案第一行的选题标题，不允许回退到文件名: {input_path}"
        )
    folder_name = sanitize_windows_name(deck_title.strip())
    return ensure_dir(archive_root / folder_name)


def resolve_cover_title(deck_title: str | None, specs: list[dict[str, Any]], input_path: pathlib.Path) -> str:
    if deck_title:
        return deck_title.strip()
    if specs:
        return (specs[0].get("deck_title") or specs[0].get("title_text") or specs[0].get("page_title") or "").strip()
    raise ValueError(f"输入文案缺少选题标题，封面标题只允许使用文案第一行的选题标题: {input_path}")


def resolve_title_folder_name(deck_title: str | None, input_path: pathlib.Path) -> str:
    if not deck_title or not deck_title.strip():
        raise ValueError(
            f"输入文案缺少选题标题，正式归档子文件夹名只允许使用文案第一行的选题标题，不允许回退到文件名: {input_path}"
        )
    return sanitize_windows_name(deck_title.strip())


def find_existing_workdir_by_title(deck_title: str, outputs_root: pathlib.Path) -> pathlib.Path | None:
    title_folder = sanitize_windows_name(deck_title)
    preferred_names = [
        f"{title_folder}（成稿）-work",
        f"{title_folder}-work",
    ]
    for name in preferred_names:
        candidate = outputs_root / name
        if candidate.exists():
            return candidate.resolve()
    if not outputs_root.exists():
        return None
    for candidate in sorted(outputs_root.iterdir()):
        if not candidate.is_dir():
            continue
        if candidate.name.startswith(title_folder) and candidate.name.endswith("-work"):
            return candidate.resolve()
    return None


def validate_existing_workdir(workdir: pathlib.Path) -> None:
    required = [
        workdir / "pages-spec.json",
        workdir / "codex-handoff.json",
        workdir / "codex-workflow.txt",
        workdir / "prompts",
    ]
    missing = [str(path) for path in required if not path.exists()]
    if missing:
        raise FileNotFoundError("现有工作包不完整，缺少以下文件或目录:\n" + "\n".join(missing))
    handoff = json.loads((workdir / "codex-handoff.json").read_text(encoding="utf-8"))
    specs_payload = json.loads((workdir / "pages-spec.json").read_text(encoding="utf-8"))
    errors: list[str] = []
    if handoff.get("identity_reference_mode") != "face-anchor-plus-action-reference":
        errors.append("顶层 identity_reference_mode 不是 face-anchor-plus-action-reference")
    if handoff.get("identity_face_anchor_image") != FACE_ANCHOR_FILENAME:
        errors.append(f"顶层 identity_face_anchor_image 不是 {FACE_ANCHOR_FILENAME}")
    pages = handoff.get("pages", [])
    if not pages:
        errors.append("codex-handoff.json 缺少 pages")
    for item in pages:
        page_index = item.get("page_index", "?")
        if item.get("role_calibration_image"):
            errors.append(f"page-{page_index} 仍包含已淘汰字段 role_calibration_image")
        if not item.get("role_face_anchor_image"):
            errors.append(f"page-{page_index} 缺少 role_face_anchor_image")
        if "role_action_reference_image" not in item:
            errors.append(f"page-{page_index} 缺少 role_action_reference_image 字段")
        if item.get("role_reference_strategy") != "anchor-plus-pose":
            errors.append(f"page-{page_index} 的 role_reference_strategy 不是 anchor-plus-pose")
    for cover in handoff.get("covers", []):
        cover_type = cover.get("cover_type", "?")
        if cover.get("role_calibration_image"):
            errors.append(f"{cover_type} 仍包含已淘汰字段 role_calibration_image")
        if not cover.get("role_face_anchor_image"):
            errors.append(f"{cover_type} 缺少 role_face_anchor_image")
        if "role_action_reference_image" not in cover:
            errors.append(f"{cover_type} 缺少 role_action_reference_image 字段")
    for cover in specs_payload.get("cover_outputs", []):
        cover_type = cover.get("cover_type", "?")
        if cover.get("role_calibration_image"):
            errors.append(f"pages-spec cover_outputs/{cover_type} 仍包含已淘汰字段 role_calibration_image")
        if not cover.get("role_face_anchor_image"):
            errors.append(f"pages-spec cover_outputs/{cover_type} 缺少 role_face_anchor_image")
    if errors:
        raise ValueError(
            "现有工作包格式已淘汰，需要用当前版本重新生成工作包，而不是继续复用旧包:\n"
            + "\n".join(f"- {error}" for error in errors)
        )


def build_archive_expected_outputs(
    selected_pages: list[int],
    archive_dir: pathlib.Path,
    include_covers: bool,
) -> list[str]:
    outputs = [str(archive_dir / f"page-{page_index:02d}.png") for page_index in selected_pages]
    if include_covers:
        outputs.extend([
            str(archive_dir / "cover-3x4.png"),
            str(archive_dir / "cover-4x3.png"),
        ])
    return outputs


def build_codex_render_prompt(
    *,
    deck_title: str,
    workdir: pathlib.Path,
    archive_dir: pathlib.Path,
    reference_image: pathlib.Path,
    handoff_path: pathlib.Path,
    specs_path: pathlib.Path,
    workflow_path: pathlib.Path,
    selected_pages: list[int],
    include_covers: bool,
) -> str:
    page_targets = ", ".join(f"page-{page_index:02d}.png" for page_index in selected_pages)
    cover_targets = "cover-3x4.png, cover-4x3.png" if include_covers else "无"
    expected_outputs = build_archive_expected_outputs(selected_pages, archive_dir, include_covers)
    expected_outputs_text = "\n".join(f"- {path}" for path in expected_outputs)
    return (
        f"你正在当前这条正式出图对话里执行当前选题的出图任务。只执行现有 work package，不重写 skill 规则，不重新生成 handoff，不走任何降级路径。\n\n"
        f"选题：{deck_title}\n"
        f"工作包目录：{workdir}\n"
        f"正式成品目录：{archive_dir}\n"
        f"唯一主脸锚点参考图：{reference_image}\n\n"
        f"必须严格执行以下硬约束：\n"
        f"1. 只读取并执行现有 work package：\n"
        f"- {handoff_path}\n"
        f"- {specs_path}\n"
        f"- {workflow_path}\n"
        f"2. 人物身份基准只允许：{reference_image.name}。左上正脸单图是唯一主脸锚点；允许配合固定人物特征文字做辅助稳定，并按页面语义额外带 1 张动作单图，但不得引入真人锚点图、姿态库身份图或第二身份源。\n"
        f"3. 只允许 Codex 出图。不得使用任何外部 API、兼容模型、live 后端、CLI 替代渲染链路、本地其他渲染器。\n"
        f"4. 必须服从 handoff 中现有约束，尤其是：allowed_render_mode = codex、codex_only_rendering_required = true、external_api_rendering_forbidden = true。\n"
        f"5. 输出必须直接写回正式归档目录，不能只停留在缓存、附件或临时目录。\n\n"
        f"执行要求：\n"
        f"1. 本次开始生成前，必须确认当前这条正式出图对话已经重新真实看到这张主脸锚点图；不是只读路径、不是只沿用上一次记忆。固定人物特征文字只允许辅助稳定，不得替代原图。如果当前对话没有真实图片输入，必须停止并明确汇报：当前执行环境无法真实挂载角色主脸锚点参考图，本次停止生成，未进入正式出图链路。\n"
        f"2. 再校验 work package 是否存在且约束字段完整。\n"
        f"3. 按 page_index 顺序依次生成正文页：{page_targets}。\n"
        f"4. 每生成一张，都必须回头对照主脸锚点图，检查是不是同一个人；重点核对眼镜、肤色、发型、脸型、年龄感、白西装黑内搭，以及发际线方向、额头高度、眼镜外轮廓、脸宽脸长比、鼻口间距是否一致。不像就只重做该页，不能带着漂移继续后续页面。\n"
        f"5. 正文完成后，生成封面：{cover_targets}。\n"
        f"6. 每生成一张，立刻写回 handoff 指定的 output_image。\n"
        f"7. 结束前逐项核对正式目录中的目标文件是否齐全，并把结果汇报清楚。\n\n"
        f"本次目标输出：\n"
        f"{expected_outputs_text}\n"
    )


def build_codex_thread_job_payload(
    *,
    deck_title: str,
    workdir: pathlib.Path,
    archive_dir: pathlib.Path,
    handoff_path: pathlib.Path,
    specs_path: pathlib.Path,
    workflow_path: pathlib.Path,
    reference_image: pathlib.Path,
    selected_pages: list[int],
    include_covers: bool,
) -> dict[str, Any]:
    expected_outputs = build_archive_expected_outputs(selected_pages, archive_dir, include_covers)
    return {
        "job_type": "codex-desktop-same-conversation-render",
        "deck_title": deck_title,
        "work_package_dir": str(workdir),
        "archive_dir": str(archive_dir),
        "handoff_file": str(handoff_path),
        "pages_spec_file": str(specs_path),
        "workflow_file": str(workflow_path),
        "formal_execution_mode": "same-conversation-single-reference-attachment",
        "reference_image": {
            "path": str(reference_image),
            "attachment_required": True,
            "must_be_real_image_input": True,
            "refresh_required_every_run": True,
            "identity_lock_mode": "face-anchor-plus-action-reference",
            "attach_once_per_new_topic_conversation": True,
        },
        "action_reference_policy": {
            "action_reference_required": True,
            "action_reference_count_max": 1,
            "identity_redefinition_forbidden": True,
        },
        "selected_pages": selected_pages,
        "include_covers": include_covers,
        "expected_outputs": expected_outputs,
        "render_rules": {
            "allowed_render_mode": "codex",
            "codex_only_rendering_required": True,
            "external_api_rendering_forbidden": True,
            "direct_writeback_required": True,
            "stop_if_reference_not_attached": True,
            "rerender_scope": "single-page-only",
            "reference_image_must_be_shown_every_run": True,
            "identity_review_required_after_each_page": True,
        },
        "archive_delivery_mode": "direct-or-copy-then-clean",
        "intermediate_render_allowed": True,
        "intermediate_images_must_be_deleted_after_archive": True,
        "archive_completion_requires_formal_dir_only": True,
        "execution_entrypoint": "current-conversation-only",
        "auto_attachment_capability_status": "blocked-no-explicit-tool-api",
        "current_conversation_rendering_required": True,
        "current_conversation_rendering_allowed_if_reference_visible": True,
        "blocked_reason_if_no_image_attachment": "当前执行环境无法真实挂载角色主脸锚点参考图，本次停止生成，未进入正式出图链路。",
        "identity_baseline_image": reference_image.name,
        "fixed_identity_traits_required": True,
        "fixed_identity_traits_mode": "auxiliary-only",
        "secondary_identity_sources_forbidden": True,
        "identity_reference_mode": "face-anchor-plus-action-reference",
        "identity_face_anchor_variant": "top-left-frontal",
        "identity_face_anchor_required": True,
        "identity_face_anchor_image": reference_image.name,
        "action_reference_required": True,
        "expression_exaggeration_limit": "reduced",
        "chibi_drift_forbidden": True,
        "identity_review_checklist": [
            "眼镜一致",
            "肤色一致",
            "发型一致",
            "脸型一致",
            "年龄感一致",
            "白西装黑内搭一致",
            "额头与发际线方向一致",
            "眼镜外轮廓一致",
            "脸宽脸长比一致",
            "五官间距一致",
            "成熟感一致",
        ],
    }


def write_codex_thread_job_files(
    *,
    deck_title: str,
    workdir: pathlib.Path,
    archive_dir: pathlib.Path,
    handoff_path: pathlib.Path,
    specs_path: pathlib.Path,
    workflow_path: pathlib.Path,
    reference_image: pathlib.Path,
    selected_pages: list[int],
    include_covers: bool,
) -> tuple[pathlib.Path, pathlib.Path]:
    prompt_path = workdir / "codex-render-thread-prompt.md"
    job_path = workdir / "codex-render-job.json"
    prompt_text = build_codex_render_prompt(
        deck_title=deck_title,
        workdir=workdir,
        archive_dir=archive_dir,
        reference_image=reference_image,
        handoff_path=handoff_path,
        specs_path=specs_path,
        workflow_path=workflow_path,
        selected_pages=selected_pages,
        include_covers=include_covers,
    )
    job_payload = build_codex_thread_job_payload(
        deck_title=deck_title,
        workdir=workdir,
        archive_dir=archive_dir,
        handoff_path=handoff_path,
        specs_path=specs_path,
        workflow_path=workflow_path,
        reference_image=reference_image,
        selected_pages=selected_pages,
        include_covers=include_covers,
    )
    write_text(prompt_path, prompt_text)
    write_json(job_path, job_payload)
    return prompt_path, job_path


def rebind_existing_work_package_outputs(
    workdir: pathlib.Path,
    archive_dir: pathlib.Path,
) -> tuple[pathlib.Path, pathlib.Path, pathlib.Path, pathlib.Path | None, pathlib.Path | None]:
    handoff_path = workdir / "codex-handoff.json"
    workflow_path = workdir / "codex-workflow.txt"
    specs_path = workdir / "pages-spec.json"

    handoff = json.loads(handoff_path.read_text(encoding="utf-8"))
    old_archive_dir = str(handoff.get("archive_dir") or "")
    old_final_archive_dir = str(handoff.get("final_archive_dir") or "")
    target_archive_dir = str(archive_dir)

    for item in handoff.get("pages", []):
        page_index = int(item.get("page_index", 0))
        if page_index > 0:
            item["output_image"] = str(archive_dir / f"page-{page_index:02d}.png")
    for cover in handoff.get("covers", []):
        cover_type = str(cover.get("cover_type") or "").strip()
        if cover_type:
            cover["output_image"] = str(archive_dir / f"{cover_type}.png")
    handoff["archive_dir"] = target_archive_dir
    handoff["final_archive_dir"] = target_archive_dir
    write_json(handoff_path, handoff)

    specs_payload = json.loads(specs_path.read_text(encoding="utf-8"))
    for cover in specs_payload.get("cover_outputs", []):
        cover_type = str(cover.get("cover_type") or "").strip()
        if cover_type:
            cover["output_image"] = str(archive_dir / f"{cover_type}.png")
    write_json(specs_path, specs_payload)

    workflow_text = workflow_path.read_text(encoding="utf-8")
    if old_archive_dir:
        workflow_text = workflow_text.replace(old_archive_dir, target_archive_dir)
    if old_final_archive_dir and old_final_archive_dir != old_archive_dir:
        workflow_text = workflow_text.replace(old_final_archive_dir, target_archive_dir)
    write_text(workflow_path, workflow_text)

    prompt_path: pathlib.Path | None = None
    job_path: pathlib.Path | None = None
    reference_image = pathlib.Path(
        next(
            (
                page.get("role_face_anchor_image")
                for page in handoff.get("pages", [])
                if page.get("role_face_anchor_image")
            ),
            "",
        )
    )
    if reference_image.exists():
        deck_title = str(
            handoff.get("deck_title")
            or next((page.get("deck_title") for page in handoff.get("pages", []) if page.get("deck_title")), "")
            or next((page.get("page_title") for page in handoff.get("pages", []) if page.get("page_title")), "")
        ).strip()
        selected_pages = [int(page["page_index"]) for page in handoff.get("pages", []) if page.get("page_index")]
        include_covers = bool(handoff.get("covers"))
        prompt_path, job_path = write_codex_thread_job_files(
            deck_title=deck_title,
            workdir=workdir,
            archive_dir=archive_dir,
            handoff_path=handoff_path,
            specs_path=specs_path,
            workflow_path=workflow_path,
            reference_image=reference_image,
            selected_pages=selected_pages,
            include_covers=include_covers,
        )
        handoff["codex_render_thread_prompt_file"] = str(prompt_path)
        handoff["codex_render_job_file"] = str(job_path)
        handoff["codex_render_thread_required"] = True
        handoff["codex_render_current_thread_allowed"] = True
        write_json(handoff_path, handoff)
    return handoff_path, workflow_path, specs_path, prompt_path, job_path


def split_sentences(text: str) -> list[str]:
    normalized = clean_input_text(text)
    if not normalized:
        return []
    sentences: list[str] = []
    for raw_line in normalized.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        parts = re.findall(r"[^。！？；!?;\n]+[。！？；!?;]?", line)
        for part in parts:
            sentence = part.strip()
            if sentence:
                sentences.append(sentence)
    return sentences


def extract_step_label(sentence: str) -> str | None:
    match = STEP_LABEL_RE.match(sentence)
    if not match:
        return None
    label = match.group("label").strip()
    return label.rstrip("，,：:。 ")


def strip_step_label(sentence: str) -> str:
    return STEP_LABEL_RE.sub("", sentence, count=1).lstrip("，,：:。 ")


def normalize_short_label(text: str) -> str:
    normalized = re.sub(r"[：:，,。！？；、“”\"'（）()《》【】\[\]]+", " ", text)
    normalized = re.sub(r"\s+", " ", normalized).strip()
    return normalized


def extract_source_terms(text: str, limit: int = 4, max_len: int = 8) -> list[str]:
    terms: list[str] = []
    for token in re.findall(r"[\u4e00-\u9fffA-Za-z0-9]+", normalize_short_label(text)):
        if len(token) < 2 or len(token) > max_len:
            continue
        if token not in terms:
            terms.append(token)
        if len(terms) >= limit:
            break
    return terms[:limit]


def chunk_body_to_title_lines(chunk: str) -> tuple[str, str, bool]:
    lines = [line.strip() for line in chunk.splitlines() if line.strip()]
    if not lines:
        return "", "", False
    bold_title_index: int | None = None
    bold_title = ""
    for index, line in enumerate(lines):
        candidate = extract_bold_title(line)
        if candidate:
            bold_title_index = index
            bold_title = candidate
            break
    if bold_title_index is not None:
        title = bold_title
        body_lines = [strip_markdown_bold(line) for idx, line in enumerate(lines) if idx != bold_title_index]
        body = "\n".join(line for line in body_lines if line).strip()
        return title, body, True
    title = strip_markdown_bold(lines[0])
    body_lines = [strip_markdown_bold(line) for line in lines[1:]]
    body = "\n".join(line for line in body_lines if line).strip()
    return title, body, False


def build_cover_support_labels(specs: list[dict[str, Any]], cover_title: str) -> list[str]:
    labels: list[str] = []
    for spec in specs[:3]:
        labels.extend(spec.get("support_labels") or [])
        labels.extend(spec.get("flow_labels") or [])
    normalized: list[str] = []
    for label in labels:
        cleaned = normalize_short_label(str(label))
        if 1 <= len(cleaned) <= 10 and cleaned not in normalized:
            normalized.append(cleaned)
    title_tokens = set(re.findall(r"[\u4e00-\u9fffA-Za-z0-9]+", cover_title))
    filtered = [label for label in normalized if label not in title_tokens]
    return filtered[:3]


def build_cover_cue_phrases(specs: list[dict[str, Any]], cover_title: str) -> list[str]:
    first_spec = specs[0] if specs else {}
    candidates = []
    for phrase in first_spec.get("cue_phrases") or []:
        cleaned = normalize_short_label(str(phrase))
        if not cleaned or cleaned == normalize_short_label(cover_title):
            continue
        if len(cleaned) > 16:
            continue
        candidates.append(cleaned)
    return unique_preserve_order(candidates)


def trim_cover_accent_blocks(
    cue_phrases: list[str],
    support_labels: list[str],
) -> tuple[list[str], list[str], int]:
    cue_candidates = unique_preserve_order(cue_phrases)
    support_candidates = [label for label in unique_preserve_order(support_labels) if label not in cue_candidates]
    selected_cue_phrases = cue_candidates[:1]
    remaining_slots = COVER_ACCENT_BLOCKS_MAX - len(selected_cue_phrases)
    selected_support_labels = support_candidates[:max(remaining_slots, 0)]
    total_blocks = len(selected_cue_phrases) + len(selected_support_labels)
    return selected_cue_phrases, selected_support_labels, total_blocks


def choose_cover_action_guidance(cover_title: str) -> str:
    if re.search(r"(赚不到|没钱|收入|变现|副业)", cover_title):
        return "角色要明确演出“白天很强、下班后变现受阻”的反差感，可以做出掏空口袋、推不开收入门、对着结果箭头发力的动作"
    if re.search(r"(误区|坑|陷阱)", cover_title):
        return "角色要像在识别和拆开误区，动作带有阻挡、拆解、指出问题的感觉"
    if re.search(r"(增长|突破|翻盘|逆袭)", cover_title):
        return "角色动作要更主动推进，带有冲刺、推开、跃迁、拉升的感觉"
    return "角色动作要服务标题核心矛盾，带有明显判断、推动、拆解或受阻后的反击感"


def choose_cover_expression_guidance(cover_title: str) -> str:
    if re.search(r"(赚不到|没钱|受困|卡住|误区)", cover_title):
        return "表情要有强烈反差，混合不甘、警醒、较真和一点点被现实卡住后的发狠感"
    if re.search(r"(机会|增长|翻盘|突破)", cover_title):
        return "表情要更兴奋、更笃定，带有抓住机会和主动突破的状态"
    return "表情要明显、夸张、可信，优先体现判断感、推进感和主题里的情绪张力"


def build_cover_outputs(
    specs: list[dict[str, Any]],
    deck_title: str | None,
    archive_dir: pathlib.Path,
    input_path: pathlib.Path,
) -> list[dict[str, Any]]:
    cover_title = resolve_cover_title(deck_title, specs, input_path)
    raw_support_labels = build_cover_support_labels(specs, cover_title)
    raw_cue_phrases = build_cover_cue_phrases(specs, cover_title)
    cue_phrases, support_labels, accent_blocks_total = trim_cover_accent_blocks(raw_cue_phrases, raw_support_labels)
    first_spec = specs[0] if specs else {}
    role_expression_tags = unique_preserve_order(first_spec.get("role_expression_tags") or [])[:4]
    role_action_tags = unique_preserve_order(first_spec.get("role_action_tags") or [])[:4]
    items: list[dict[str, Any]] = []
    for cover in COVER_OUTPUT_SPECS:
        items.append({
            "cover_type": cover["cover_type"],
            "aspect_ratio": cover["aspect_ratio"],
            "title_text": cover_title,
            "display_title": cover_title,
            "background_color": COVER_BG,
            "skin_tone_base": SKIN_TONE_BASE,
            "line_color": ODD_LINE,
            "role_required": True,
            "style_mode": "hand-drawn-visual-note",
            "support_labels_allowed": True,
            "support_labels": support_labels,
            "cue_phrases": cue_phrases,
            "source_text_only_visual_mode": True,
            "inferred_visuals_forbidden": True,
            "role_count_min": 1,
            "role_count_max": 3,
            "role_usage_mode": "supporting-only",
            "role_density_preference": "medium-support",
            "cover_accent_blocks_min": COVER_ACCENT_BLOCKS_MIN,
            "cover_accent_blocks_max": COVER_ACCENT_BLOCKS_MAX,
            "cover_accent_blocks_total": accent_blocks_total,
            "role_action_tags": role_action_tags,
            "role_expression_tags": role_expression_tags,
            "action_guidance": choose_cover_action_guidance(cover_title),
            "expression_guidance": choose_cover_expression_guidance(cover_title),
            "composition_guidance": cover["composition_guidance"],
            "title_zone": cover["title_zone"],
            "output_image": str(archive_dir / cover["filename"]),
        })
    return items


def build_cover_reference_chain(
    cover_ref_source: dict[str, Any],
    role: dict[str, Any],
) -> tuple[dict[str, Any] | None, dict[str, Any] | None, str]:
    face_anchor = select_face_anchor_profile(role)
    action_ref, action_reason = select_action_reference_profile(cover_ref_source, role)
    if not action_ref:
        action_ref = find_reference_profile_by_filename(role, "ip-ref-09-idea.png")
        action_reason = "cover-default"
    return face_anchor, action_ref, action_reason


def attach_cover_role_chain(
    cover_outputs: list[dict[str, Any]],
    specs: list[dict[str, Any]],
    role: dict[str, Any],
) -> None:
    cover_ref_source = specs[0] if specs else {}
    face_anchor_ref, action_ref, action_reason = build_cover_reference_chain(cover_ref_source, role)
    for cover_spec in cover_outputs:
        cover_spec["role_face_anchor_image"] = str(face_anchor_ref["file_path"]) if face_anchor_ref else ""
        cover_spec["role_action_reference_image"] = str(action_ref["file_path"]) if action_ref else ""
        cover_spec["role_action_reference_reason"] = action_reason or (action_ref.get("name", "") if action_ref else "")


def build_reference_lock_metadata() -> dict[str, Any]:
    return {
        "reference_lock_required": True,
        "text_only_render_forbidden": True,
        "renderer_must_attach_reference_images": True,
        "reference_image_refresh_required_every_run": True,
        "reference_image_must_be_shown_to_model_every_run": True,
        "identity_reference_mode": "face-anchor-plus-action-reference",
        "identity_baseline_image": FACE_ANCHOR_FILENAME,
        "fixed_identity_traits_required": True,
        "fixed_identity_traits_mode": "auxiliary-only",
        "secondary_identity_sources_forbidden": True,
        "identity_face_anchor_variant": "top-left-frontal",
        "identity_face_anchor_required": True,
        "identity_face_anchor_image": FACE_ANCHOR_FILENAME,
        "action_reference_required": True,
        "expression_exaggeration_limit": "reduced",
        "chibi_drift_forbidden": True,
        "required_reference_inputs": ["face-anchor", "action-reference?"],
        "optional_non_identity_references": [],
        "fail_if_reference_images_not_actually_attached": True,
        "codex_only_rendering_required": True,
        "external_api_rendering_forbidden": True,
        "allowed_render_mode": "codex",
        "identity_review_required_after_each_page": True,
        "identity_review_reference_source": FACE_ANCHOR_FILENAME,
        "identity_review_checklist": [
            "黑框眼镜一致",
            "暖肤色一致",
            "背头侧分黑发一致",
            "白西装黑内搭一致",
            "年轻亚洲男性一致",
            "手绘线稿风一致",
            "发际线方向一致",
            "额头高度一致",
            "眼镜外轮廓一致",
            "脸宽脸长比一致",
            "鼻口间距一致",
            "年龄感一致",
        ],
        "archive_delivery_mode": "direct-or-copy-then-clean",
        "intermediate_render_allowed": True,
        "intermediate_images_must_be_deleted_after_archive": True,
        "archive_completion_requires_formal_dir_only": True,
        "source_text_only_visual_mode": True,
        "inferred_visuals_forbidden": True,
        "role_count_min": 1,
        "role_count_max": 3,
        "role_usage_mode": "supporting-only",
        "role_density_preference": "supporting-first",
    }


def parse_explicit_pages(text: str) -> list[dict[str, Any]]:
    chunks = [chunk.strip() for chunk in re.split(r"(?m)^\s*---\s*$", text) if chunk.strip()]
    pages: list[dict[str, Any]] = []
    for idx, chunk in enumerate(chunks, start=1):
        title, body, title_from_bold = chunk_body_to_title_lines(chunk)
        if not title:
            continue
        pages.append({
            "index": idx,
            "title": title,
            "body": body,
            "title_from_bold": title_from_bold,
            "pagination_mode": "explicit",
            "section_kind": "manual",
            "step_label": "",
            "step_index": None,
            "step_page_index": 1,
            "step_page_total": 1,
            "is_continued": False,
        })
    if not pages:
        raise ValueError("没有解析到任何页面，请检查文案中是否有内容")
    return pages


def group_sentences_by_structure(text: str) -> list[dict[str, Any]]:
    sentences = split_sentences(text)
    if not sentences:
        return []
    blocks: list[dict[str, Any]] = []
    current_intro: list[str] = []
    current_step: dict[str, Any] | None = None

    for sentence in sentences:
        label = extract_step_label(sentence)
        if label:
            if current_step:
                blocks.append(current_step)
            elif current_intro:
                blocks.append({"kind": "intro", "sentences": current_intro[:]})
                current_intro = []
            current_step = {
                "kind": "step",
                "step_label": label,
                "sentences": [sentence],
            }
            continue
        if current_step:
            current_step["sentences"].append(sentence)
        else:
            current_intro.append(sentence)

    if current_step:
        blocks.append(current_step)
    elif current_intro:
        blocks.append({"kind": "intro", "sentences": current_intro[:]})

    if blocks and blocks[-1]["kind"] == "step":
        trailing_outro = peel_outro_sentences(blocks[-1]["sentences"])
        if trailing_outro:
            blocks[-1]["sentences"] = blocks[-1]["sentences"][:-len(trailing_outro)]
            blocks.append({"kind": "outro", "sentences": trailing_outro})
    return [block for block in blocks if block.get("sentences")]


def peel_outro_sentences(sentences: list[str]) -> list[str]:
    if len(sentences) < 2:
        return []
    outro: list[str] = []
    for sentence in reversed(sentences):
        if SUMMARY_SENTENCE_RE.search(sentence):
            outro.insert(0, sentence)
        elif outro:
            break
        else:
            break
    if len(outro) == 1:
        return outro
    return outro


def count_chars(sentences: list[str]) -> int:
    return sum(len(sentence) for sentence in sentences)


def pick_split_index(sentences: list[str], min_left: int = 2, min_right: int = 2) -> int:
    if len(sentences) <= min_left + min_right:
        return max(min_left, len(sentences) - min_right)
    candidate_indices = []
    for idx in range(min_left, len(sentences) - min_right + 1):
        left = sentences[idx - 1]
        right = sentences[idx]
        if PAGE_BREAK_RE.search(right) or PAGE_BREAK_RE.search(left):
            candidate_indices.append(idx)
    if candidate_indices:
        midpoint = len(sentences) / 2.0
        return min(candidate_indices, key=lambda idx: abs(idx - midpoint))
    return max(min_left, len(sentences) // 2)


def summarize_page_title(step_label: str, sentences: list[str]) -> str:
    if not sentences:
        return ""
    text = strip_step_label(sentences[0])
    text = re.sub(r"^(那|所以|其实|就是|然后|你要知道|你现在最该做的，是|你现在最该做的|你该补的是|也就是说|正确顺序是|很多人|普通人最怕的不是|今天分享一个很扎心的真相，?)", "", text)
    text = re.sub(r"[，。！？；,;:：].*$", "", text).strip()
    text = text[:18].strip()
    if len(text) < 4:
        text = strip_step_label("".join(sentences))[:12].strip()
    if not text:
        return ""
    return text


def summarize_intro_title(sentences: list[str], page_index: int) -> str:
    if not sentences:
        return f"导语{page_index}"
    first = sentences[0]
    title = re.sub(r"[，。！？；,;:：].*$", "", first).strip()
    title = re.sub(r"^(今天分享一个|今天我就把|先说一个|先讲一个)", "", title).strip()
    if len(title) < 4:
        title = first[:12].strip()
    return title[:18] or f"导语{page_index}"


def summarize_hook(sentences: list[str], fallback: str = "") -> str:
    if not sentences:
        return fallback[:18]
    text = strip_step_label(sentences[0])
    text = re.sub(r"^(今天分享一个|今天我就把|所以|其实|也就是说|真正厉害的人，不是|真正厉害的人是)", "", text).strip()
    text = re.sub(r"[。！？；!?;].*$", "", text).strip()
    text = text[:22].strip()
    return text or fallback[:18]


def strip_leading_spoken_filler(text: str) -> str:
    return re.sub(
        r"^(对，你没听错。?|你没听错。?|好，?|那这时候你一定会问[:：]?|记住，?|所以，?|但是，?|因为，?|其实，?|也就是说，?)",
        "",
        text.strip(),
    ).strip()


def is_explicit_step_title(title: str) -> bool:
    return bool(EXPLICIT_STEP_TITLE_RE.search(title.strip()))


def normalize_step_title_entry(entry: str) -> str:
    cleaned = entry.strip().strip("。！？；; ")
    marker_match = EXPLICIT_STEP_TITLE_RE.search(cleaned)
    if not marker_match:
        return cleaned
    marker = marker_match.group(0).strip()
    rest = cleaned[marker_match.end():].lstrip("：:，,、 ")
    rest = re.sub(r"^(就是|是)\s*", "", rest).strip()
    if not rest:
        return marker
    return f"{marker}：{rest.strip().strip('。！？；; ')}"


def extract_step_title_layers(title: str, body: str) -> list[str]:
    layers: list[str] = []
    sources = [title] + [line.strip() for line in body.splitlines() if line.strip()]
    for idx, source in enumerate(sources):
        for match in STEP_TITLE_ENTRY_RE.finditer(source):
            entry = normalize_step_title_entry(match.group("entry"))
            if entry and "：" not in entry and idx + 1 < len(sources):
                next_line = strip_leading_spoken_filler(sources[idx + 1]).strip().strip("。！？；; ")
                if next_line and not is_explicit_step_title(next_line):
                    entry = f"{entry}：{next_line}"
            if entry and entry not in layers:
                layers.append(entry)
            if len(layers) >= 2:
                return layers[:2]
    return layers[:2]


def is_complete_sentence_title(text: str) -> bool:
    normalized = strip_leading_spoken_filler(strip_step_label(text)).strip()
    if len(normalized) < 8:
        return False
    if normalized.endswith(("，", ",", "：", ":", "；", ";", "、")):
        return False
    return any(mark in normalized for mark in ["，", "。", "：", "不是", "而是", "就是", "才是", "必须", "不要", "记住"])


def score_full_hook_title(text: str) -> int:
    normalized = strip_leading_spoken_filler(strip_step_label(text)).strip()
    score = 0
    if 10 <= len(normalized) <= 24:
        score += 2
    elif len(normalized) >= 8:
        score += 1
    for marker in ["不是", "而是", "就是", "才是", "记住", "真正", "最重要", "答案", "路线图", "闭环", "长出来"]:
        if marker in normalized:
            score += 2
    if normalized.startswith(("这时候", "然后", "那这时候", "好，这就", "对，你没听错")):
        score -= 2
    return score


def extract_full_hook_title(candidates: list[str], fallback: str = "") -> str:
    scored_complete_candidates: list[tuple[int, str]] = []
    for raw in candidates:
        cleaned = strip_leading_spoken_filler(strip_step_label(raw))
        cleaned = cleaned.replace("“", "").replace("”", "").strip()
        cleaned = re.sub(r"[。！？；!?;]+\s*$", "", cleaned).strip()
        if not cleaned:
            continue
        if len(cleaned) > 30:
            cleaned = cleaned[:30].rstrip("，,：:；; ")
        if is_complete_sentence_title(cleaned):
            scored_complete_candidates.append((score_full_hook_title(cleaned), cleaned))
    if scored_complete_candidates:
        scored_complete_candidates.sort(key=lambda item: item[0], reverse=True)
        return scored_complete_candidates[0][1]
    for raw in candidates:
        cleaned = strip_leading_spoken_filler(strip_step_label(raw))
        cleaned = re.sub(r"[。！？；!?;]+\s*$", "", cleaned).strip()
        if len(cleaned) >= 8:
            return cleaned[:30].rstrip("，,：:；; ")
    fallback_clean = strip_leading_spoken_filler(strip_step_label(fallback)).strip()
    return fallback_clean[:30].rstrip("，,：:；; ") or fallback[:30].rstrip("，,：:；; ")


def compress_sentence_to_phrase(sentence: str, limit: int = 14) -> str:
    text = strip_step_label(sentence)
    text = re.sub(r"^(今天分享一个|今天我就把|很多人|所以|但是|因为|其实|也就是说|正确顺序是|真正该补的是|你现在最该做的，是|你现在最该做的|普通人最怕的不是)", "", text).strip()
    text = re.sub(r"[。！？；!?;].*$", "", text).strip()
    text = text.replace("“", "").replace("”", "")
    hard_splits = ["，", "；", "：", "、", "不是", "而是", "先", "再", "最后"]
    for splitter in hard_splits:
        if splitter in text and len(text) > limit:
            left = text.split(splitter, 1)[0].strip()
            if 2 <= len(left) <= limit:
                text = left
                break
    return text[:limit].strip("，,：:；; ")


def extract_cue_phrases(sentences: list[str], limit_points: int = 4) -> list[str]:
    points: list[str] = []
    for sentence in sentences:
        phrase = compress_sentence_to_phrase(sentence)
        if not phrase:
            continue
        if phrase not in points:
            points.append(phrase)
        if len(points) >= limit_points:
            break
    return points[:limit_points]


def extract_visual_keywords(title: str, body: str, icons: list[str], page_type: str) -> list[str]:
    del icons, page_type
    keywords: list[str] = []
    for source in [title] + [line.strip() for line in body.splitlines() if line.strip()][:3]:
        for term in extract_source_terms(source, limit=2, max_len=8):
            if term not in keywords:
                keywords.append(term)
            if len(keywords) >= 4:
                break
        if len(keywords) >= 4:
            break
    return keywords[:4]


def extract_visual_action_cues(title: str, body: str, page_type: str) -> list[str]:
    haystack = f"{title}\n{body}"
    cues: list[str] = []
    mapping = [
        (r"(打篮球|篮球)", "打篮球"),
        (r"(打游戏|游戏)", "打游戏"),
        (r"(写作|写东西|写内容)", "写作"),
        (r"(电脑|笔记本|办公)", "用电脑"),
        (r"(销售|成交|客户|沟通)", "谈销售"),
        (r"(对比|不是|而是|区别|反差)", "对比"),
        (r"(流程|步骤|推进|执行|行动)", "推进"),
        (r"(讲|解释|说明|拆解)", "讲解"),
        (r"(惊讶|扎心|戳中|反差)", "惊讶"),
        (r"(困惑|拖延|卡住|误区|没结果)", "困惑"),
        (r"(总结|收尾|结论|确认)", "确认"),
        (r"(目标|验证|深度工作)", "深度工作"),
    ]
    for pattern, cue in mapping:
        if re.search(pattern, haystack) and cue not in cues:
            cues.append(cue)
    fallback = {
        "对比页": ["对比", "讲解"],
        "流程页": ["推进", "讲解"],
        "观点页": ["讲解", "强调"],
        "总结页": ["确认", "讲解"],
        "步骤页": ["讲解"],
    }
    for cue in fallback.get(page_type, ["讲解"]):
        if cue not in cues:
            cues.append(cue)
    return cues[:4]


def extract_visual_scene_cues(title: str, body: str, page_type: str) -> list[str]:
    del page_type
    haystack = f"{title}\n{body}"
    cues: list[str] = []
    explicit_terms = [
        "步骤",
        "流程",
        "对比",
        "服务",
        "客户",
        "系统",
        "关系",
        "问题",
        "方法",
        "过程",
        "节点",
        "筛选",
        "展示",
    ]
    for term in explicit_terms:
        if term in haystack and term not in cues:
            cues.append(term)
    return cues[:4]


def extract_flow_labels(page: dict[str, Any], cue_phrases: list[str], text_density_mode: str) -> list[str]:
    labels: list[str] = []
    if page.get("step_label"):
        labels.append(str(page["step_label"]))
    if text_density_mode == "graph-first":
        return labels[:1]
    limit = 1 if text_density_mode == "text-light" else 2
    for phrase in cue_phrases:
        if phrase not in labels:
            labels.append(phrase)
        if len(labels) >= limit:
            break
    return labels[:limit]


def is_long_title(title: str) -> bool:
    plain = strip_step_label(title).replace("“", "").replace("”", "").strip()
    return len(plain) >= 24


def choose_display_title_mode(title: str, page_type: str) -> str:
    if is_long_title(title):
        return "compressed"
    if page_type in {"对比页", "流程页"} and len(strip_step_label(title)) >= 20:
        return "compressed"
    return "exact"


def choose_title_strategy(page: dict[str, Any], page_type: str, deck_title: str | None) -> str:
    if int(page.get("index", 0) or 0) == 1 and deck_title and not page.get("title_from_bold"):
        return "deck-title-exact"
    step_layers = extract_step_title_layers(str(page.get("title") or ""), str(page.get("body") or ""))
    if len(step_layers) >= 2:
        return "explicit-step-with-subtitle"
    if len(step_layers) == 1:
        return "explicit-step-title"
    return "extracted-full-hook"


def extract_guided_title_prefix(title: str) -> str:
    stripped = title.strip()
    match = GUIDE_TITLE_PREFIX_RE.match(stripped)
    if match:
        return match.group("prefix").strip()
    return extract_step_label(stripped) or ""


def strip_guided_title_prefix(title: str) -> str:
    stripped = title.strip()
    prefix = extract_guided_title_prefix(stripped)
    if not prefix:
        return stripped
    rest = stripped[len(prefix):].lstrip("，,：:。 、|｜-")
    return rest.strip()


def build_display_title(title: str, page_type: str, mode: str) -> str:
    if mode == "exact":
        return title
    prefix = extract_guided_title_prefix(title)
    base_title = strip_guided_title_prefix(title) if prefix else title
    limit = 16 if page_type in {"对比页", "总结页"} else 18
    if prefix:
        remainder_limit = max(6, limit - len(prefix) - 1)
        compressed = compress_sentence_to_phrase(base_title, limit=remainder_limit)
        if compressed:
            return f"{prefix}｜{compressed}"
        if base_title:
            return f"{prefix}｜{base_title[:remainder_limit]}"
        return prefix
    compressed = compress_sentence_to_phrase(base_title, limit=limit)
    return compressed or base_title[:limit]


def extract_subtitle_notes(title: str, mode: str, limit: int = 2) -> list[str]:
    if mode != "compressed":
        return []
    text = strip_guided_title_prefix(title).replace("“", "").replace("”", "").strip()
    chunks = re.split(r"[，。！？；!?;]", text)
    notes: list[str] = []
    for chunk in chunks:
        chunk = chunk.strip()
        if not chunk:
            continue
        phrase = compress_sentence_to_phrase(chunk, limit=14)
        if phrase and phrase not in notes:
            notes.append(phrase)
        if len(notes) >= limit:
            break
    return notes[:limit]


def choose_text_density_mode(index: int, page_type: str, title: str, scene_concept: str, section_kind: str) -> str:
    del title, scene_concept
    if index == 1 or section_kind == "intro":
        return "text-light"
    if page_type in {"总结页"} or section_kind == "outro":
        return "text-light"
    if page_type in {"步骤页", "观点页", "对比页", "流程页"} or section_kind == "step":
        return "graph-first"
    return "text-light"


def extract_support_labels(
    page_type: str,
    scene_concept: str,
    cue_phrases: list[str],
    flow_labels: list[str],
    visual_scene_cues: list[str],
    visual_keywords: list[str],
    title: str,
    body: str,
    text_density_mode: str,
) -> list[str]:
    del page_type, scene_concept
    if text_density_mode == "minimal":
        return []
    max_labels = 3 if text_density_mode == "text-light" else 5
    labels: list[str] = []
    body_lines = [line.strip() for line in body.splitlines() if line.strip()][:4]
    phrase_sources = list(visual_scene_cues) + list(visual_keywords) + list(cue_phrases) + list(flow_labels) + [title] + body_lines
    for source in phrase_sources:
        if len(labels) >= max_labels:
            break
        for term in extract_source_terms(str(source), limit=3, max_len=6):
            if term not in labels:
                labels.append(term)
            if len(labels) >= max_labels:
                break
    if not labels:
        fallback_phrase = compress_sentence_to_phrase(str(title), limit=6)
        if fallback_phrase:
            labels.append(fallback_phrase)
    return labels[:max_labels]


def derive_role_action_tags(page_type: str, visual_action_cues: list[str], role_action: str, role_pose_hint: str) -> list[str]:
    tags: list[str] = []
    for source in visual_action_cues + [role_action, role_pose_hint]:
        for label in ["讲解", "对比", "推进", "强调", "惊讶", "困惑", "确认", "打篮球", "打游戏", "写作", "用电脑", "谈销售", "深度工作"]:
            if label in source and label not in tags:
                tags.append(label)
    page_type_map = {
        "对比页": "对比",
        "流程页": "推进",
        "总结页": "确认",
        "观点页": "强调",
    }
    mapped = page_type_map.get(page_type)
    if mapped and mapped not in tags:
        tags.append(mapped)
    return tags[:5]


def derive_role_expression_tags(page_type: str, role_expression: str, title: str, body: str) -> list[str]:
    tags: list[str] = []
    haystack = "\n".join([page_type, title, body, role_expression])
    keyword_map = [
        ("惊讶", "惊讶"),
        ("质疑", "质疑"),
        ("困惑", "困惑"),
        ("思考", "思考"),
        ("顿悟", "顿悟"),
        ("判断", "判断"),
        ("解释", "解释"),
        ("推进", "推进"),
        ("执行", "执行"),
        ("确认", "确认"),
        ("号召", "号召"),
        ("笃定", "笃定"),
        ("警觉", "警觉"),
        ("专注", "专注"),
        ("鼓舞", "鼓舞"),
    ]
    for keyword, tag in keyword_map:
        if keyword in haystack and tag not in tags:
            tags.append(tag)
    fallback = {
        "对比页": ["判断", "解释"],
        "流程页": ["推进", "执行"],
        "总结页": ["确认", "号召"],
        "观点页": ["思考", "强调"],
    }
    for tag in fallback.get(page_type, ["说明"]):
        if tag not in tags:
            tags.append(tag)
    return tags[:4]


def build_role_variation_candidate(
    action_family: str,
    expression_family: str,
    action: str,
    expression: str,
    pose_hint: str,
    action_tags: list[str],
    expression_tags: list[str],
    preferred_framing: str | None = None,
    preferred_position: str | None = None,
) -> dict[str, Any]:
    return {
        "action_family": action_family,
        "expression_family": expression_family,
        "role_action": action,
        "role_expression": expression,
        "role_pose_hint": pose_hint,
        "role_action_tags": unique_preserve_order([action_family] + action_tags),
        "role_expression_tags": unique_preserve_order([expression_family] + expression_tags),
        "preferred_framing": preferred_framing,
        "preferred_position": preferred_position,
    }


def extend_role_variation_candidates(candidates: list[dict[str, Any]], additions: list[dict[str, Any]]) -> list[dict[str, Any]]:
    existing_keys = {
        (
            item["action_family"],
            item["expression_family"],
            item["role_action"],
            item["role_expression"],
        )
        for item in candidates
    }
    for item in additions:
        key = (
            item["action_family"],
            item["expression_family"],
            item["role_action"],
            item["role_expression"],
        )
        if key not in existing_keys:
            candidates.append(item)
            existing_keys.add(key)
    return candidates


def select_role_reference_profiles(spec: dict[str, Any], role: dict[str, Any], limit: int = 2) -> list[dict[str, Any]]:
    profiles = role.get("reference_profiles") or []
    if not profiles:
        return []
    desired = set((spec.get("role_action_tags") or []) + (spec.get("visual_action_cues") or []) + (spec.get("visual_scene_cues") or []))
    preferred_variant = "exaggerated"
    page_type = spec.get("page_type", "")
    if page_type in {"总结页"}:
        preferred_variant = "exaggerated"
    if spec.get("section_kind") in {"intro"} and spec.get("page_index") not in {1}:
        preferred_variant = "standard"
    if spec.get("role_scale") == "small-accent" and page_type not in {"对比页", "观点页"}:
        preferred_variant = "standard"
    scored: list[tuple[int, int, int, dict[str, Any]]] = []
    for item in profiles:
        tags = set(item.get("tags") or []) | set(item.get("use_cases") or [])
        score = len(desired & tags)
        variant_bonus = 1 if item.get("variant") == preferred_variant else 0
        intensity_bonus = 1 if item.get("expression_intensity") in {"high", "extreme-controlled"} else 0
        if score > 0 or item.get("is_base_pose"):
            scored.append((score, variant_bonus, intensity_bonus, item))
    if not scored:
        return profiles[:limit]
    scored.sort(key=lambda pair: (-pair[0], -pair[1], -pair[2], pair[3].get("name", "")))
    selected: list[dict[str, Any]] = []
    calibration = next(
        (
            item
            for _, _, _, item in scored
            if "cartoon-style" in set(item.get("tags") or [])
            or "hair-calibration" in set(item.get("tags") or [])
            or "肤色校准" in set(item.get("use_cases") or [])
        ),
        None,
    )
    if calibration:
        selected.append(calibration)
    for _, _, _, item in scored:
        if item in selected:
            continue
        selected.append(item)
        if len(selected) >= limit:
            break
    return selected[:limit]


def select_role_calibration_profile(role: dict[str, Any], preferred_variant: str = "standard") -> dict[str, Any] | None:
    profiles = role.get("reference_profiles") or []
    candidates = [
        item for item in profiles
        if "cartoon-style" in set(item.get("tags") or [])
        or "hair-calibration" in set(item.get("tags") or [])
        or "肤色校准" in set(item.get("use_cases") or [])
    ]
    if not candidates:
        return None
    exact = next((item for item in candidates if item.get("variant") == preferred_variant), None)
    if exact:
        return exact
    return candidates[0]


def find_reference_profile_by_filename(role: dict[str, Any], filename: str) -> dict[str, Any] | None:
    profiles = role.get("reference_profiles") or []
    return next((item for item in profiles if pathlib.Path(item.get("path", "")).name == filename), None)


def select_face_anchor_profile(role: dict[str, Any]) -> dict[str, Any] | None:
    anchor = find_reference_profile_by_filename(role, FACE_ANCHOR_FILENAME)
    if anchor:
        return anchor
    return select_role_calibration_profile(role, preferred_variant="standard")


def select_action_reference_profile(spec: dict[str, Any], role: dict[str, Any]) -> tuple[dict[str, Any] | None, str]:
    haystack = " ".join(
        str(value)
        for value in [
            spec.get("page_type"),
            spec.get("scene_concept"),
            spec.get("page_title"),
            " ".join(spec.get("role_action_tags") or []),
            " ".join(spec.get("role_expression_tags") or []),
            " ".join(spec.get("visual_action_cues") or []),
            " ".join(spec.get("visual_scene_cues") or []),
        ]
        if value
    )
    for keywords, filenames in ACTION_REFERENCE_BY_INTENT:
        matched = next((keyword for keyword in keywords if keyword in haystack), None)
        if not matched:
            continue
        for filename in filenames:
            profile = find_reference_profile_by_filename(role, filename)
            if profile:
                return profile, matched
    return None, ""


def split_role_reference_chain(
    spec: dict[str, Any],
    role: dict[str, Any],
    preferred_variant: str = "standard",
) -> tuple[dict[str, Any] | None, dict[str, Any] | None, str]:
    face_anchor = select_face_anchor_profile(role)
    action_ref, action_reason = select_action_reference_profile(spec, role)
    return face_anchor, action_ref, action_reason


def choose_scene_concept(title: str, page_type: str, section_kind: str, index: int) -> str:
    if is_cover_like_page(index, title):
        return "cover-hero"
    if page_type == "流程页":
        return "path-flow"
    if page_type == "对比页":
        return "contrast-split"
    if page_type == "总结页":
        return "summary-stage"
    return "concept-note-scene"


def choose_scene_layout_type(scene_concept: str, role_scale: str, index: int) -> str:
    if scene_concept == "cover-hero":
        return "cover-hero"
    mapping = {
        "pit-map": "pit-map",
        "production-line": "production-line",
        "path-flow": "path-flow",
        "contrast-split": "contrast-split",
        "blocked-structure": "blocked-structure",
        "summary-stage": "summary-stage",
        "cover-hero": "cover-hero",
    }
    return mapping.get(scene_concept, "concept-board")


def choose_title_layout_mode(scene_concept: str, index: int, role_scale: str) -> str:
    if scene_concept == "cover-hero":
        return "center-hero"
    if index == 1:
        return "top-span" if role_scale != "large-focus" else "left-anchor"
    if scene_concept in {"pit-map", "path-flow"}:
        return "top-span"
    if scene_concept == "contrast-split":
        return "left-anchor" if role_scale != "large-focus" else "right-anchor"
    if scene_concept == "production-line":
        return "left-anchor"
    if scene_concept == "blocked-structure":
        return "right-anchor"
    if scene_concept == "summary-stage":
        return "center-hero"
    return "inline-scene"


def choose_title_alignment(title_layout_mode: str) -> str:
    mapping = {
        "center-hero": "center",
        "top-span": "center",
        "left-anchor": "left",
        "right-anchor": "right",
        "inline-scene": "left",
    }
    return mapping.get(title_layout_mode, "left")


def choose_title_anchor_zone(scene_concept: str, title_layout_mode: str) -> str:
    if title_layout_mode == "center-hero":
        return "center" if scene_concept == "cover-hero" else "top"
    if title_layout_mode == "top-span":
        return "top"
    mapping = {
        "top-span": "top",
        "left-anchor": "upper-left",
        "right-anchor": "upper-right",
        "inline-scene": "inline",
    }
    return mapping.get(title_layout_mode, "top")


def choose_title_flow_direction(scene_concept: str, title_layout_mode: str, title: str) -> str:
    if title_layout_mode == "top-span":
        return "horizontal"
    if title_layout_mode == "center-hero" and len(title) > 16:
        return "stacked"
    if scene_concept in {"path-flow", "blocked-structure"}:
        return "scene-following"
    return "horizontal"


def is_cover_like_page(index: int, title: str) -> bool:
    if index != 1:
        return False
    return bool(re.search(r"(封面|cover|目录|总览|导读|开场|首页)", title, re.I))


def choose_title_layout_guidance(
    scene_concept: str,
    title_layout_mode: str,
    title_alignment: str,
    title_anchor_zone: str,
    title_flow_direction: str,
) -> str:
    if title_layout_mode == "center-hero":
        return "标题作为中心锤点或上中锤点，与人物和主场景一起形成海报式主视觉"
    if title_layout_mode == "top-span":
        return "标题横跨顶部或偏上区域，下面完整展开结构和路径"
    if title_layout_mode == "left-anchor":
        return "标题偏左作为文字锚点，右侧或中部承接人物与主结构"
    if title_layout_mode == "right-anchor":
        return "标题偏右作为文字锚点，左侧或中部承接人物与主结构"
    return (
        f"标题融入场景内部，不像固定条幅；对齐方式 {title_alignment}，锚点区域 {title_anchor_zone}，"
        f"标题走向 {title_flow_direction}。人物不要默认站成画面中心，优先服从场景结构。"
    )


def choose_role_scene_relationship(scene_concept: str, role_scale: str, index: int) -> str:
    if index == 1:
        return "主题优先，第一张也必须先服从原文结构，角色只能作为辅助演示元素进入主结构"
    if role_scale == "small-accent":
        return "主题和主场景优先，角色缩小进入结构内部，作为辅助行动者参与表达"
    if role_scale == "large-focus":
        return "仅在封面或极少数强主视觉页允许更强存在感，但仍要服务原文"
    if scene_concept in {"pit-map", "production-line", "path-flow"}:
        return "主题和结构优先，角色必须走进结构、操作节点或被结构限制，不能站在结构外解说"
    return "主题优先，角色必须与主场景产生直接动作关系，主场景优先于角色单体"


def choose_role_in_scene_mode(title: str, page_type: str, scene_concept: str, index: int) -> str:
    haystack = f"{title}\n{scene_concept}"
    if index == 1:
        return "supporting-embedded"
    if page_type == "对比页":
        return "supporting-bridge"
    if page_type == "流程页" or scene_concept in {"production-line", "path-flow"}:
        return "supporting-operator"
    if re.search(r"(误区|踩坑|坑|卡住|地图|总览|路线图)", haystack):
        return "supporting-structure-interactor"
    if page_type == "总结页":
        return "supporting-summary-actor"
    return "supporting-structure-interactor"


def choose_role_in_scene_guidance(role_in_scene_mode: str, page_type: str, scene_concept: str) -> str:
    if role_in_scene_mode == "supporting-bridge":
        return "人物必须站进对比结构之间，同时连接左右两侧信息，像桥接、切换或权衡，但只能做辅助。"
    if role_in_scene_mode == "supporting-operator":
        return "人物必须进入流程、路径或操作结构内部，沿节点推进、推动、跨越或操作主机制，但不能取代主体内容。"
    if role_in_scene_mode == "supporting-summary-actor":
        return "人物必须在主场景内部承担动作或结论收束，不做画外主持人，不悬浮在空白边缘。"
    return "人物必须走进结构内部，触碰、踩入、跨过、被包围或被结构限制，不能独立站在主结构外。"


def choose_supporting_micro_visuals(icon_set: list[str], scene_concept: str) -> list[str]:
    items: list[str] = []
    for icon in icon_set[:3]:
        if icon not in items:
            items.append(icon)
    scene_defaults = {
        "pit-map": ["箭头", "标签框", "警示线"],
        "production-line": ["流程箭头", "节点框", "连接线"],
        "path-flow": ["箭头", "分叉点", "连接线"],
        "contrast-split": ["左右分栏", "对照线", "连接箭头"],
        "blocked-structure": ["阻断线", "门槛线", "框线"],
        "summary-stage": ["收束箭头", "重点框", "确认线"],
        "cover-hero": ["标题框", "强调线", "分栏提示"],
    }
    for extra in scene_defaults.get(scene_concept, []):
        if extra not in items:
            items.append(extra)
    return items[:6]


def get_density_profile(scene_concept: str, section_kind: str, is_continued: bool, index: int) -> dict[str, int]:
    if index == 1 or scene_concept == "cover-hero":
        return {"talking_points": 2, "micro_visuals": 4, "visual_keywords": 4, "icons": 2, "support_labels": 3, "primary_blocks": 2, "secondary_blocks": 0}
    if is_continued:
        return {"talking_points": 1, "micro_visuals": 4, "visual_keywords": 4, "icons": 2, "support_labels": 3, "primary_blocks": 2, "secondary_blocks": 0}
    if scene_concept in {"pit-map", "production-line", "path-flow", "blocked-structure"}:
        return {"talking_points": 1, "micro_visuals": 4, "visual_keywords": 4, "icons": 2, "support_labels": 5, "primary_blocks": 2, "secondary_blocks": 0}
    if scene_concept == "summary-stage":
        return {"talking_points": 2, "micro_visuals": 4, "visual_keywords": 4, "icons": 2, "support_labels": 4, "primary_blocks": 2, "secondary_blocks": 1}
    if section_kind == "step":
        return {"talking_points": 1, "micro_visuals": 4, "visual_keywords": 4, "icons": 2, "support_labels": 4, "primary_blocks": 2, "secondary_blocks": 0}
    return {"talking_points": 2, "micro_visuals": 4, "visual_keywords": 4, "icons": 2, "support_labels": 3, "primary_blocks": 2, "secondary_blocks": 0}


def extract_following_step_labels(pages: list[dict[str, str]], current_index: int, limit: int) -> list[str]:
    labels: list[str] = []
    for page in pages[current_index:]:
        title = str(page.get("title") or "").strip()
        if not title:
            continue
        label = strip_guided_title_prefix(title)
        label = compress_sentence_to_phrase(label or title, limit=6)
        if label and label not in labels:
            labels.append(label)
        if len(labels) >= limit:
            break
    return labels[:limit]


def build_page_dict(
    index: int,
    title: str,
    sentences: list[str],
    pagination_mode: str,
    section_kind: str,
    step_label: str = "",
    step_index: int | None = None,
    step_page_index: int = 1,
    step_page_total: int = 1,
) -> dict[str, Any]:
    return {
        "index": index,
        "title": title,
        "body": "\n".join(sentences).strip(),
        "pagination_mode": pagination_mode,
        "section_kind": section_kind,
        "step_label": step_label,
        "step_index": step_index,
        "step_page_index": step_page_index,
        "step_page_total": step_page_total,
        "is_continued": step_page_index > 1,
    }


def split_intro_block(sentences: list[str]) -> list[list[str]]:
    if len(sentences) <= 4 and count_chars(sentences) <= 140:
        return [sentences]
    if len(sentences) <= 2:
        return [sentences]
    split_idx = None
    for idx in range(1, len(sentences)):
        if INTRO_BREAK_RE.search(sentences[idx]):
            split_idx = idx
            break
    if split_idx is None:
        split_idx = pick_split_index(sentences, min_left=1, min_right=1)
    chunks = [sentences[:split_idx], sentences[split_idx:]]
    return [chunk for chunk in chunks if chunk]


def split_step_block(sentences: list[str]) -> list[list[str]]:
    if len(sentences) <= 5 and count_chars(sentences) <= 180:
        return [sentences]
    chunks: list[list[str]] = []
    remaining = sentences[:]
    while remaining:
        if len(remaining) <= 5 and count_chars(remaining) <= 180:
            chunks.append(remaining)
            break
        take = remaining[:6]
        split_idx = pick_split_index(take, min_left=2, min_right=2 if len(take) >= 4 else 1)
        chunk = remaining[:split_idx]
        chunks.append(chunk)
        remaining = remaining[split_idx:]
    return [chunk for chunk in chunks if chunk]


def parse_auto_pages(text: str) -> list[dict[str, Any]]:
    blocks = group_sentences_by_structure(text)
    if not blocks:
        raise ValueError("没有解析到任何页面，请检查文案中是否有内容")

    pages: list[dict[str, Any]] = []
    page_index = 1
    step_counter = 0
    for block in blocks:
        kind = block["kind"]
        sentences = block["sentences"]
        if kind == "intro":
            intro_chunks = split_intro_block(sentences)
            for idx, chunk in enumerate(intro_chunks, start=1):
                pages.append(build_page_dict(
                    index=page_index,
                    title=summarize_intro_title(chunk, idx),
                    sentences=chunk,
                    pagination_mode="auto",
                    section_kind="intro",
                ))
                page_index += 1
            continue
        if kind == "step":
            step_counter += 1
            step_chunks = split_step_block(sentences)
            step_label = block["step_label"]
            total = len(step_chunks)
            for idx, chunk in enumerate(step_chunks, start=1):
                title = step_label if idx == 1 else summarize_page_title(step_label, chunk)
                pages.append(build_page_dict(
                    index=page_index,
                    title=title,
                    sentences=chunk,
                    pagination_mode="auto",
                    section_kind="step",
                    step_label=step_label,
                    step_index=step_counter,
                    step_page_index=idx,
                    step_page_total=total,
                ))
                page_index += 1
            continue
        if kind == "outro":
            pages.append(build_page_dict(
                index=page_index,
                title="总结",
                sentences=sentences,
                pagination_mode="auto",
                section_kind="outro",
            ))
            page_index += 1
    return pages


def parse_pages(text: str) -> list[dict[str, Any]]:
    if re.search(r"(?m)^\s*---\s*$", text):
        return parse_explicit_pages(text)
    return parse_auto_pages(text)


def parse_page_selection(selection: str | None, total_pages: int) -> list[int]:
    if not selection:
        return list(range(1, total_pages + 1))
    picked: set[int] = set()
    for chunk in selection.split(","):
        part = chunk.strip()
        if not part:
            continue
        match = PAGE_REF_RE.match(part)
        if not match:
            raise ValueError(f"无法解析页码范围: {part}")
        start = int(match.group(1))
        end = int(match.group(2) or start)
        if start > end:
            start, end = end, start
        if start < 1 or end > total_pages:
            raise ValueError(f"页码超出范围: {part}，当前总页数 {total_pages}")
        picked.update(range(start, end + 1))
    return sorted(picked)


def infer_page_type(title: str, body: str, index: int, total: int) -> str:
    haystack = f"{title}\n{body}"
    if index == total and re.search(r"(总结|结尾|最后|收尾|行动|号召)", haystack):
        return "总结页"
    if re.search(r"(对比|不是|而是|vs|区别|差异)", haystack, re.I):
        return "对比页"
    if re.search(r"(步骤|第一步|第二步|第三步|流程|路径|方法|先|再|然后)", haystack):
        return "流程页"
    if re.search(r"(认知|观点|本质|为什么|核心|真相|底层)", haystack):
        return "观点页"
    if index == total:
        return "总结页"
    return "步骤页"


def build_role_variation_candidates(spec: dict[str, Any]) -> list[dict[str, Any]]:
    page_type = spec["page_type"]
    title = spec["page_title"]
    body = spec.get("page_body", "")
    scene_concept = spec.get("scene_concept", "")
    haystack = f"{title}\n{body}"
    candidates: list[dict[str, Any]] = []

    if re.search(r"(误区|踩坑|坑)", haystack):
        extend_role_variation_candidates(candidates, [
            build_role_variation_candidate(
                "跨越",
                "警觉",
                "全身跨过坑位、障碍或错误节点，边躲边指认问题来源",
                "高夸张的警觉感和惊讶感，像刚发现关键风险",
                "让人物真正踩进或跨过误区结构，身体被坑位、箭头或警示标签包围",
                ["跨越", "误区", "推进"],
                ["警觉", "惊讶", "判断"],
                preferred_framing="full-body",
                preferred_position="inside-structure",
            ),
            build_role_variation_candidate(
                "阻挡",
                "质疑",
                "伸手挡住错误路径或标签，像在当场否决一个误区",
                "高夸张的质疑感和否定感，眉眼明显收紧",
                "用阻挡、拦截或推开的演法处理误区，不要只做普通解释",
                ["阻挡", "误区", "强调"],
                ["质疑", "判断", "警觉"],
                preferred_framing="three-quarter",
                preferred_position="mid-structure",
            ),
        ])
    if re.search(r"(路径|流程|步骤|顺序|验证)", haystack):
        extend_role_variation_candidates(candidates, [
            build_role_variation_candidate(
                "推进",
                "专注",
                "沿着流程节点推进、连接或逐步操作，把动作落在路径结构上",
                "高夸张的专注感和执行感，像在带着页面往前走",
                "动作必须贴着箭头、节点或流程线推进，不允许站在流程外讲解",
                ["推进", "流程", "操作"],
                ["专注", "执行", "推进"],
                preferred_framing="full-body",
                preferred_position="mid-path",
            ),
            build_role_variation_candidate(
                "连接",
                "带动",
                "跨步越过前后节点，伸手把两个阶段直接连起来",
                "高夸张的带动感和节奏感，像在把流程硬推起来",
                "用跨步、拉动、连接前后节点的方式表现顺序，不要只指着流程说",
                ["连接", "跨越", "流程"],
                ["带动", "推进", "执行"],
                preferred_framing="full-body",
                preferred_position="mid-path",
            ),
        ])
    if re.search(r"(赚钱|收入|变现|销售|成交|客户)", haystack):
        extend_role_variation_candidates(candidates, [
            build_role_variation_candidate(
                "对比",
                "判断",
                "一手对照投入，一手对照结果或收入符号，做强反差判断",
                "高夸张的判断感和解释感，像在拆开投入产出关系",
                "让人物同时连到投入和结果两边，形成明显的反差结构",
                ["对比", "收入", "谈销售"],
                ["判断", "解释", "质疑"],
                preferred_framing="three-quarter",
                preferred_position="center-bridge",
            ),
            build_role_variation_candidate(
                "确认",
                "困惑",
                "看向空空的钱包、订单或结果面板，做出落差确认",
                "高夸张的困惑感和落差感，像发现努力没有变成结果",
                "动作里要带结果落空的对象，不要只做口头说明姿势",
                ["确认", "收入", "结果"],
                ["困惑", "惊讶", "判断"],
                preferred_framing="three-quarter",
                preferred_position="inside-right",
            ),
        ])
    if re.search(r"(写作|内容|文案|输出)", haystack):
        extend_role_variation_candidates(candidates, [
            build_role_variation_candidate(
                "写作",
                "专注",
                "低头手写、标注或敲字，把内容动作直接放进页面结构里",
                "高夸张的专注感和思考感，像正在现场产出内容",
                "让人物直接操作纸张、键盘或文案节点，不要脱离内容结构",
                ["写作", "操作", "拆解"],
                ["专注", "思考", "执行"],
                preferred_framing="three-quarter",
                preferred_position="inside-left",
            ),
        ])
    if re.search(r"(电脑|系统|工具|模型)", haystack):
        extend_role_variation_candidates(candidates, [
            build_role_variation_candidate(
                "操作",
                "专注",
                "直接操作电脑、面板或系统节点，像在现场验证工具链路",
                "高夸张的专注感和判断感，眼神盯住屏幕或面板反馈",
                "让人物和工具结构发生操作关系，不要退回普通主持手势",
                ["操作", "用电脑", "推进"],
                ["专注", "判断", "执行"],
                preferred_framing="three-quarter",
                preferred_position="inside-right",
            ),
        ])

    page_type_defaults = {
        "对比页": [
            build_role_variation_candidate(
                "对比",
                "判断",
                "双手分别指向两侧内容，做取舍、对照和判断",
                "高夸张的判断感和解释感，眉眼和手势都更明显",
                "站进对比结构中间，让双手和视线同时连到两侧内容",
                ["对比", "判断", "讲解"],
                ["判断", "解释", "强调"],
                preferred_framing="three-quarter",
                preferred_position="center-bridge",
            ),
            build_role_variation_candidate(
                "桥接",
                "解释",
                "站在中间把两侧结构连接起来，像在搭桥解释差异",
                "高夸张的解释感和带动感，像把两组内容当场串起来",
                "身体要压在左右结构之间，做桥接动作，不要离场站桩",
                ["桥接", "对比", "连接"],
                ["解释", "带动", "判断"],
                preferred_framing="full-body",
                preferred_position="center-bridge",
            ),
            build_role_variation_candidate(
                "取舍",
                "质疑",
                "一手推开错误侧，一手拉近正确侧，明确做取舍",
                "高夸张的质疑感和判断感，像在当场否掉一个方向",
                "动作必须落在两侧结构上，形成明显推开和拉近的方向差",
                ["取舍", "对比", "强调"],
                ["质疑", "判断", "警觉"],
                preferred_framing="three-quarter",
                preferred_position="center-bridge",
            ),
        ],
        "流程页": [
            build_role_variation_candidate(
                "推进",
                "执行",
                "沿着流程节点逐步推进，手势落在关键连接点上",
                "高夸张的推进感和执行感，动作更有带动性",
                "不要站在流程旁边讲，必须沿着节点真正推进",
                ["推进", "流程", "执行"],
                ["执行", "推进", "专注"],
                preferred_framing="full-body",
                preferred_position="mid-path",
            ),
            build_role_variation_candidate(
                "操作",
                "专注",
                "蹲进或靠近某个关键节点，像在当场操作流程开关",
                "高夸张的专注感和判断感，注意力锁在结构细节上",
                "让人物缩进流程内部做操作，不要做飘在外面的讲解动作",
                ["操作", "流程", "拆解"],
                ["专注", "判断", "执行"],
                preferred_framing="three-quarter",
                preferred_position="mid-structure",
            ),
            build_role_variation_candidate(
                "连接",
                "带动",
                "伸手连接前后节点或把节点往前拖动，形成递进感",
                "高夸张的带动感和推进感，像在给流程加速度",
                "人物要和前后节点都发生连接，不要只指其中一个点",
                ["连接", "推进", "流程"],
                ["带动", "推进", "专注"],
                preferred_framing="full-body",
                preferred_position="mid-path",
            ),
        ],
        "总结页": [
            build_role_variation_candidate(
                "收束",
                "确认",
                "面向观众做收束确认，像在把重点压成一个结论",
                "高夸张的确认感和号召感，像在收束全场",
                "动作要稳住画面，把散开的信息收回到结论上",
                ["收束", "确认", "总结"],
                ["确认", "号召", "笃定"],
                preferred_framing="three-quarter",
                preferred_position="inside-center",
            ),
            build_role_variation_candidate(
                "展开",
                "鼓舞",
                "双手向外展开，把最终结论和行动口径一起打开",
                "高夸张的鼓舞感和笃定感，像在给最后一下推动",
                "总结页可以稳，但不要变成静止站姿，要有开场或号召手势",
                ["展开", "总结", "号召"],
                ["鼓舞", "笃定", "确认"],
                preferred_framing="half-body",
                preferred_position="inside-right",
            ),
            build_role_variation_candidate(
                "落点",
                "笃定",
                "伸手按住结论框、终点节点或总结标签，像在落锤",
                "高夸张的笃定感和确认感，像把结论直接钉住",
                "把动作压到最终落点上，不要只站着看总结框",
                ["落点", "确认", "总结"],
                ["笃定", "确认", "判断"],
                preferred_framing="three-quarter",
                preferred_position="inside-left",
            ),
        ],
        "观点页": [
            build_role_variation_candidate(
                "强调",
                "思考",
                "指向重点词，边强调边做思考或拆解动作",
                "高夸张的思考感和解释感，表情更戏剧化",
                "把手势和视线落在重点词或核心结构上，不要离题发挥",
                ["强调", "讲解", "拆解"],
                ["思考", "解释", "判断"],
                preferred_framing="three-quarter",
                preferred_position="inside-left",
            ),
            build_role_variation_candidate(
                "质疑",
                "质疑",
                "挑眉、半摊手，像在反问一个默认认知",
                "高夸张的质疑感和困惑感，形成明显的认知冲突",
                "人物要像在质问结构里的某个假设，不要变成普通说明员",
                ["质疑", "强调", "判断"],
                ["质疑", "困惑", "思考"],
                preferred_framing="half-body",
                preferred_position="inside-right",
            ),
            build_role_variation_candidate(
                "顿悟",
                "顿悟",
                "突然前探或握拳，像刚抓到一个关键认知",
                "高夸张的顿悟感和兴奋感，眼神更亮、动作更直接",
                "让人物像在主结构里突然点亮关键节点，而不是只讲观点",
                ["顿悟", "强调", "推进"],
                ["顿悟", "惊讶", "鼓舞"],
                preferred_framing="half-body",
                preferred_position="inside-center",
            ),
        ],
        "步骤页": [
            build_role_variation_candidate(
                "演示",
                "说明",
                "围绕页面核心概念做演示和拆解",
                "高夸张但稳定可辨认的说明型表情",
                "角色结合标题内容演示，不要站桩，不要只是普通主持姿态",
                ["演示", "讲解", "拆解"],
                ["说明", "解释", "思考"],
                preferred_framing="three-quarter",
                preferred_position="inside-right",
            ),
            build_role_variation_candidate(
                "拆解",
                "思考",
                "一边拆开结构，一边盯住某个关键点做解释",
                "高夸张的思考感和判断感，像在现场拆题",
                "动作要落在结构内部，不要只对着观众解释",
                ["拆解", "讲解", "判断"],
                ["思考", "判断", "解释"],
                preferred_framing="three-quarter",
                preferred_position="inside-left",
            ),
            build_role_variation_candidate(
                "强调",
                "确认",
                "抬手压住一个重点标签，像在把这一页的主张钉牢",
                "高夸张的确认感和强调感，像在对一个关键点定调",
                "重点是压住页面关键词，不要回到通用主持姿势",
                ["强调", "确认", "讲解"],
                ["确认", "判断", "解释"],
                preferred_framing="half-body",
                preferred_position="inside-center",
            ),
        ],
    }
    extend_role_variation_candidates(candidates, page_type_defaults.get(page_type, page_type_defaults["步骤页"]))
    if scene_concept in {"pit-map", "blocked-structure"}:
        extend_role_variation_candidates(candidates, [
            build_role_variation_candidate(
                "受限",
                "困惑",
                "让身体被结构卡住、挤压或围住，像被问题机制限制住",
                "高夸张的困惑感和警觉感，像刚意识到限制条件",
                "让人物真被结构限制住，而不是站在结构外讨论它",
                ["受限", "结构", "误区"],
                ["困惑", "警觉", "质疑"],
                preferred_framing="full-body",
                preferred_position="inside-structure",
            ),
        ])
    return candidates


def build_role_position_options(spec: dict[str, Any], preferred_position: str | None = None) -> list[str]:
    page_type = spec["page_type"]
    scene_concept = spec.get("scene_concept", "")
    role_scale = spec.get("role_scale", "medium-support")
    options = [preferred_position or "", spec.get("role_position", "")]
    if page_type == "对比页":
        options.extend(["center-bridge", "inside-left", "inside-right"])
    elif page_type == "流程页" or scene_concept in {"production-line", "path-flow"}:
        options.extend(["mid-path", "mid-structure", "inside-left", "inside-right"])
    elif scene_concept in {"pit-map", "blocked-structure"}:
        options.extend(["inside-structure", "mid-structure", "inside-left", "inside-right"])
    elif role_scale == "large-focus":
        options.extend(["inside-left", "inside-right", "inside-center"])
    else:
        options.extend(["inside-left", "inside-right", "inside-center", "mid-structure"])
    return unique_preserve_order(options)


def build_role_framing_options(spec: dict[str, Any], preferred_framing: str | None = None) -> list[str]:
    page_type = spec["page_type"]
    role_scale = spec.get("role_scale", "medium-support")
    title = spec.get("page_title", "")
    options = [preferred_framing or "", spec.get("role_framing", "")]
    if re.search(r"(误区|踩坑|跨越|路径|流程|步骤|走路)", title):
        options.extend(["full-body", "three-quarter"])
    elif role_scale == "large-focus":
        options.extend(["three-quarter", "half-body"])
    elif page_type == "总结页":
        options.extend(["three-quarter", "half-body", "full-body"])
    elif role_scale == "small-accent":
        options.extend(["full-body", "three-quarter"])
    else:
        options.extend(["three-quarter", "half-body", "full-body"])
    return unique_preserve_order(options)


def choose_role_variation_staging(
    spec: dict[str, Any],
    candidate: dict[str, Any],
    recent_specs: list[dict[str, Any]],
) -> tuple[str, str, str]:
    if not recent_specs:
        return spec.get("role_position", ""), spec.get("role_framing", ""), "keep-default-staging"
    position_options = build_role_position_options(spec, candidate.get("preferred_position"))
    framing_options = build_role_framing_options(spec, candidate.get("preferred_framing"))
    previous = recent_specs[-1] if recent_specs else {}
    previous_combo = (
        previous.get("role_position"),
        previous.get("role_framing"),
        previous.get("role_composition_mode"),
    )
    best_combo = (
        spec.get("role_position", ""),
        spec.get("role_framing", ""),
        spec.get("role_composition_mode", ""),
    )
    best_score: int | None = None
    for framing in framing_options:
        for position in position_options:
            combo = (position, framing, spec.get("role_composition_mode", ""))
            score = 0
            if combo == previous_combo:
                score += 20
            if previous and position == previous.get("role_position"):
                score += 4
            if previous and framing == previous.get("role_framing"):
                score += 4
            if position == spec.get("role_position"):
                score += 1
            if framing == spec.get("role_framing"):
                score += 1
            if best_score is None or score < best_score:
                best_score = score
                best_combo = combo
    position, framing, _ = best_combo
    guard = "keep-default-staging"
    if position != spec.get("role_position") or framing != spec.get("role_framing"):
        guard = "adjust-staging-to-break-repeat"
    return position, framing, guard


def apply_role_variation_across_specs(specs: list[dict[str, Any]]) -> list[dict[str, Any]]:
    if not specs:
        return specs
    action_counts: Counter[str] = Counter()
    expression_counts: Counter[str] = Counter()
    signature_counts: Counter[str] = Counter()
    varied_specs: list[dict[str, Any]] = []
    for spec in specs:
        candidates = build_role_variation_candidates(spec)
        previous = varied_specs[-1] if varied_specs else {}
        previous2 = varied_specs[-2] if len(varied_specs) > 1 else {}
        best_candidate = candidates[0]
        best_score: int | None = None
        for idx, candidate in enumerate(candidates):
            action_family = candidate["action_family"]
            expression_family = candidate["expression_family"]
            signature = f"{action_family}|{expression_family}"
            score = idx * 4
            score += action_counts[action_family] * 5
            score += expression_counts[expression_family] * 4
            score += signature_counts[signature] * 12
            if previous:
                if action_family in (previous.get("role_action_tags") or [])[:1]:
                    score += 18
                if expression_family in (previous.get("role_expression_tags") or [])[:1]:
                    score += 18
                if candidate.get("preferred_position") == previous.get("role_position"):
                    score += 2
                if candidate.get("preferred_framing") == previous.get("role_framing"):
                    score += 2
            if previous2:
                if action_family in (previous2.get("role_action_tags") or [])[:1]:
                    score += 6
                if expression_family in (previous2.get("role_expression_tags") or [])[:1]:
                    score += 6
            if best_score is None or score < best_score:
                best_score = score
                best_candidate = candidate

        role_position, role_framing, staging_guard = choose_role_variation_staging(spec, best_candidate, varied_specs)
        updated = dict(spec)
        updated["role_action"] = best_candidate["role_action"]
        updated["role_expression"] = best_candidate["role_expression"]
        updated["role_action_tags"] = best_candidate["role_action_tags"]
        updated["role_expression_tags"] = best_candidate["role_expression_tags"]
        updated["role_pose_hint"] = f"{best_candidate['role_pose_hint']}；同一选题内避免与相邻页重复动作和表情演法"
        updated["role_position"] = role_position
        updated["role_framing"] = role_framing
        guard_parts = []
        reason_parts = []
        if best_candidate is not candidates[0]:
            guard_parts.append("switch-action-expression-family")
            reason_parts.append(
                f"为避开同题重复，本页改用“{best_candidate['action_family']} + {best_candidate['expression_family']}”组合"
            )
        if staging_guard != "keep-default-staging":
            guard_parts.append(staging_guard)
            reason_parts.append("同时调整人物景别或站位，避免连续页出现近似构图")
        if not guard_parts:
            guard_parts.append("semantic-match-primary")
            reason_parts.append("当前页首选动作已经与近邻页拉开差异，保持该页语义优先")
        updated["role_variation_guard"] = " + ".join(guard_parts)
        updated["role_variation_reason"] = "；".join(reason_parts)
        varied_specs.append(updated)
        action_counts[best_candidate["action_family"]] += 1
        expression_counts[best_candidate["expression_family"]] += 1
        signature_counts[f"{best_candidate['action_family']}|{best_candidate['expression_family']}"] += 1
    return varied_specs


def choose_role_action(page_type: str, title: str, body: str) -> str:
    if page_type == "对比页":
        return "用双手分别指向两侧内容，做对比和取舍"
    if page_type == "流程页":
        return "沿着流程节点推进、连接或跨越"
    if page_type == "总结页":
        return "面向观众做收束确认，像在总结重点"
    if page_type == "观点页":
        return "指向重点词、做思考或顿悟动作"
    return "围绕页面核心概念做演示和拆解"


def choose_role_expression(page_type: str, title: str, body: str) -> str:
    if page_type == "对比页":
        return "高夸张的判断感和解释感，眉眼和手势都更明显"
    if page_type == "流程页":
        return "高夸张的推进感和执行感，动作更有带动性"
    if page_type == "总结页":
        return "高夸张的确认感和号召感，像在收束全场"
    if page_type == "观点页":
        return "高夸张的思考、质疑或顿悟感，表情更戏剧化"
    return "高夸张但稳定可辨认的说明型表情"


def choose_icon_set(title: str, body: str, page_type: str) -> list[str]:
    haystack = f"{title}\n{body}"
    icons: list[str] = []
    keyword_map = [
        ("系统", "齿轮"),
        ("步骤", "箭头"),
        ("流程", "流程箭头"),
        ("执行", "勾选框"),
        ("认知", "灯泡"),
        ("对比", "天平"),
        ("时间", "时钟"),
        ("增长", "上升曲线"),
        ("用户", "人物群像"),
    ]
    for keyword, icon in keyword_map:
        if keyword in haystack and icon not in icons:
            icons.append(icon)
    if not icons:
        fallback = extract_source_terms(haystack, limit=2, max_len=4)
        icons.extend(fallback or ["箭头"])
    return icons[:4]


def choose_role_scale(page_type: str, section_kind: str, title: str, visual_keywords: list[str], index: int) -> str:
    if is_cover_like_page(index, title):
        return "large-focus"
    del visual_keywords
    if index == 1:
        return "medium-support"
    if page_type == "总结页":
        return "medium-support"
    if page_type in {"流程页", "对比页"}:
        return "small-accent"
    if section_kind in {"intro", "manual"}:
        return "small-accent"
    return "small-accent"


def choose_role_position(index: int, role_scale: str, section_kind: str, page_type: str, scene_concept: str) -> str:
    if page_type == "对比页":
        return "center-bridge"
    if page_type == "流程页" or scene_concept in {"production-line", "path-flow"}:
        return "mid-path"
    if page_type == "总结页":
        return "inside-center"
    if section_kind == "outro":
        return ["inside-right", "inside-left", "inside-center"][(index - 2) % 3]
    return ["inside-left", "inside-right", "inside-center", "mid-structure"][(index - 2) % 4]


def choose_role_framing(role_scale: str, page_type: str, title: str, index: int) -> str:
    haystack = title
    if is_cover_like_page(index, title) and role_scale == "large-focus":
        return "half-body"
    if role_scale == "small-accent":
        return "full-body"
    if page_type in {"总结页", "对比页", "流程页"}:
        return "three-quarter"
    return "three-quarter"


def choose_role_composition_mode(
    title: str,
    page_type: str,
    role_scale: str,
    role_position: str,
    role_framing: str,
    index: int,
) -> str:
    haystack = title
    if is_cover_like_page(index, title):
        return "cover-centered-focus"
    if index == 1:
        return "supporting-opening-scene"
    if re.search(r"(误区|踩坑|坑)", haystack):
        return "structure-first-supporting-role"
    if re.search(r"(路径|流程|步骤|验证|顺序)", haystack):
        return "supporting-operator-inside-flow"
    if page_type == "对比页":
        return "supporting-contrast-bridge"
    if role_scale == "small-accent":
        return "structure-first-small-support"
    if role_scale == "large-focus" and role_position == "center" and role_framing == "half-body":
        return "cover-centered-focus"
    return "supporting-role-in-scene"


def choose_camera_energy(page_type: str, title: str, role_scale: str) -> str:
    haystack = title
    if role_scale == "large-focus" or re.search(r"(扎心|真相|赚不到|误区|踩坑|拖延)", haystack):
        return "high"
    return "medium"


def choose_role_pose_hint(title: str, page_type: str, role_scale: str) -> str:
    haystack = title
    if is_cover_like_page(1, title):
        return "封面可以更强，但正文页先搭原文结构，再把一到三个角色嵌进去；人物只是辅助演示元素"
    if re.search(r"(7个误区|误区|踩坑|坑)", haystack):
        return "角色要走进结构里，但只能辅助指向、对照或操作，不允许用一排人物替代内容"
    if re.search(r"(赚不到|收入|变现)", haystack):
        return "做出明显反差，但不要画收入、金币、火箭等原文没说的结果物"
    if re.search(r"(路径|流程|步骤|顺序|验证)", haystack):
        return "沿着箭头、节点或路径真正推进、跨越或操作结构，动作要和路径发生直接接触"
    if page_type == "对比页":
        return "身体和手势同时连到左右两侧结构，站在中间做桥接和判断，不允许独立站在空白侧边"
    if page_type == "总结页":
        return "站稳画面，手势展开，像在做最后确认和收束"
    if role_scale == "large-focus":
        return "人物作为主视觉主体，只用于封面或极少数特殊页，正文页默认不用大主角"
    return "角色结合标题内容演示，不要站桩，不要只是普通主持姿态"


def role_scale_ratio(role_scale: str) -> str:
    return {
        "small-accent": "roughly 10% to 18% of the canvas",
        "medium-support": "roughly 20% to 35% of the canvas",
        "large-focus": "roughly 38% to 55% of the canvas",
    }.get(role_scale, "roughly 20% to 35% of the canvas")


def build_page_spec(page: dict[str, str], total: int, pages: list[dict[str, str]] | None = None) -> dict[str, Any]:
    index = int(page["index"])
    page_type = infer_page_type(page["title"], page["body"], index, total)
    odd = index % 2 == 1
    sentences = [line.strip() for line in page["body"].splitlines() if line.strip()]
    full_icon_set = choose_icon_set(page["title"], page["body"], page_type)
    page_hook = page.get("page_hook") or summarize_hook(sentences, fallback=page["title"])
    speaker_guides = sentences[: min(2, len(sentences))]
    raw_visual_keywords = page.get("visual_keywords") or extract_visual_keywords(page["title"], page["body"], full_icon_set, page_type)
    role_scale = choose_role_scale(page_type, page.get("section_kind", "manual"), page["title"], raw_visual_keywords, index)
    scene_concept = choose_scene_concept(page["title"], page_type, page.get("section_kind", "manual"), index)
    density = get_density_profile(scene_concept, page.get("section_kind", "manual"), bool(page.get("is_continued", False)), index)
    text_density_mode = page.get("text_density_mode") or choose_text_density_mode(
        index,
        page_type,
        page["title"],
        scene_concept,
        page.get("section_kind", "manual"),
    )
    icon_set = full_icon_set[: density["icons"]]
    title_strategy = page.get("title_strategy") or choose_title_strategy(page, page_type, page.get("deck_title"))
    step_title_layers = extract_step_title_layers(page["title"], page["body"])
    title_candidates = [page["title"]] + sentences
    effective_title = page["title"]
    display_title_mode = page.get("display_title_mode") or choose_display_title_mode(page["title"], page_type)
    explicit_subtitle_notes: list[str] = []
    if title_strategy == "deck-title-exact" and page.get("deck_title"):
        effective_title = str(page["deck_title"]).strip()
        display_title_mode = "exact"
    elif title_strategy == "explicit-step-with-subtitle":
        effective_title = step_title_layers[0] if step_title_layers else page["title"]
        explicit_subtitle_notes = step_title_layers[1:2]
        display_title_mode = "exact"
    elif title_strategy == "explicit-step-title":
        effective_title = step_title_layers[0] if step_title_layers else page["title"]
        display_title_mode = "exact"
    else:
        extracted_title_candidates = (
            [page["title"]] if is_complete_sentence_title(page["title"]) else []
        ) + sentences + [page["title"]]
        effective_title = extract_full_hook_title(extracted_title_candidates, fallback=page["title"])
        display_title_mode = "exact"
    display_title = page.get("display_title") or (effective_title if display_title_mode == "exact" else build_display_title(effective_title, page_type, display_title_mode))
    subtitle_notes = page.get("subtitle_notes") or explicit_subtitle_notes or extract_subtitle_notes(effective_title, display_title_mode)
    cue_phrases = (
        page.get("cue_phrases")
        or page.get("talking_points")
        or extract_cue_phrases(sentences, limit_points=density["talking_points"])
    )
    if not cue_phrases:
        cue_phrases = extract_cue_phrases([page["title"]], limit_points=max(1, density["talking_points"]))
    visual_keywords = (page.get("visual_keywords") or raw_visual_keywords)[: density["visual_keywords"]]
    scene_layout_type = choose_scene_layout_type(scene_concept, role_scale, index)
    title_layout_mode = choose_title_layout_mode(scene_concept, index, role_scale)
    title_alignment = choose_title_alignment(title_layout_mode)
    title_anchor_zone = choose_title_anchor_zone(scene_concept, title_layout_mode)
    title_flow_direction = choose_title_flow_direction(scene_concept, title_layout_mode, page["title"])
    role_position = choose_role_position(index, role_scale, page.get("section_kind", "manual"), page_type, scene_concept)
    role_framing = choose_role_framing(role_scale, page_type, page["title"], index)
    role_composition_mode = choose_role_composition_mode(page["title"], page_type, role_scale, role_position, role_framing, index)
    role_in_scene_mode = choose_role_in_scene_mode(page["title"], page_type, scene_concept, index)
    camera_energy = choose_camera_energy(page_type, page["title"], role_scale)
    role_pose_hint = choose_role_pose_hint(page["title"], page_type, role_scale)
    visual_action_cues = page.get("visual_action_cues") or extract_visual_action_cues(page["title"], page["body"], page_type)
    visual_scene_cues = page.get("visual_scene_cues") or extract_visual_scene_cues(page["title"], page["body"], page_type)
    flow_labels = page.get("flow_labels") or extract_flow_labels(page, cue_phrases, text_density_mode)
    support_labels = page.get("support_labels") or extract_support_labels(
        page_type,
        scene_concept,
        cue_phrases,
        flow_labels,
        visual_scene_cues,
        visual_keywords,
        page["title"],
        page["body"],
        text_density_mode,
    )
    support_labels = support_labels[: density.get("support_labels", 2)]
    if (
        pages
        and len(support_labels) < density.get("support_labels", 2)
        and ("误区" in page["title"] or "步骤" in page["title"])
    ):
        following_labels = extract_following_step_labels(pages, index, density.get("support_labels", 2))
        for label in following_labels:
            if label not in support_labels:
                support_labels.append(label)
            if len(support_labels) >= density.get("support_labels", 2):
                break
    primary_seed = [display_title] + subtitle_notes[:1] + cue_phrases[:1]
    primary_info_blocks = unique_preserve_order(primary_seed)[: density.get("primary_blocks", 2)]
    secondary_info_blocks: list[str] = []
    if density.get("secondary_blocks", 0) > 0 and text_density_mode != "graph-first":
        secondary_info_blocks = unique_preserve_order(
            [label for label in support_labels + flow_labels if label not in primary_info_blocks]
        )[: density.get("secondary_blocks", 1)]
    role_action = choose_role_action(page_type, page["title"], page["body"])
    role_expression = choose_role_expression(page_type, page["title"], page["body"])
    role_action_tags = page.get("role_action_tags") or derive_role_action_tags(page_type, visual_action_cues, role_action, role_pose_hint)
    role_expression_tags = page.get("role_expression_tags") or derive_role_expression_tags(page_type, role_expression, page["title"], page["body"])
    supporting_micro_visuals = choose_supporting_micro_visuals(icon_set, scene_concept)[: density["micro_visuals"]]
    return {
        "page_index": index,
        "page_title": page["title"],
        "step_title_layers": step_title_layers,
        "title_text": effective_title,
        "title_strategy": title_strategy,
        "title_extraction_required": title_strategy == "extracted-full-hook",
        "forbid_fragment_title": title_strategy == "extracted-full-hook",
        "display_title_mode": display_title_mode,
        "display_title": display_title,
        "source_title_full": page["title"],
        "subtitle_notes": subtitle_notes,
        "page_body": page["body"],
        "deck_title": page.get("deck_title"),
        "page_hook": page_hook,
        "cue_phrases": cue_phrases,
        "talking_points": cue_phrases,
        "flow_labels": flow_labels,
        "support_labels": support_labels,
        "primary_info_blocks": primary_info_blocks,
        "secondary_info_blocks": secondary_info_blocks,
        "primary_info_block_visual_required": True,
        "secondary_info_block_visual_optional": True,
        "visual_keywords": visual_keywords,
        "visual_action_cues": visual_action_cues,
        "visual_scene_cues": visual_scene_cues,
        "role_action_tags": role_action_tags,
        "role_expression_tags": role_expression_tags,
        "speaker_guides": speaker_guides,
        "density_profile": density,
        "text_density_mode": text_density_mode,
        "graph_first_mode": text_density_mode == "graph-first",
        "text_light_mode": text_density_mode == "text-light",
        "page_type": page_type,
        "page_background_mode": "odd" if odd else "even",
        "background_color": ODD_BG if odd else EVEN_BG,
        "skin_tone_base": SKIN_TONE_BASE,
        "line_color": ODD_LINE if odd else EVEN_LINE,
        "red_accent_required": True,
        "red_accent_count_min": 2,
        "red_accent_count_max": 4,
        "red_accent_allowed_targets": ["keyword", "arrow", "tag", "result", "warning"],
        "scene_priority": "scene-first",
        "scene_concept": scene_concept,
        "scene_layout_type": scene_layout_type,
        "title_layout_mode": title_layout_mode,
        "title_alignment": title_alignment,
        "title_anchor_zone": title_anchor_zone,
        "title_flow_direction": title_flow_direction,
        "title_layout_guidance": choose_title_layout_guidance(
            scene_concept,
            title_layout_mode,
            title_alignment,
            title_anchor_zone,
            title_flow_direction,
        ),
        "supporting_micro_visuals": supporting_micro_visuals,
        "source_text_only_visual_mode": True,
        "inferred_visuals_forbidden": True,
        "role_required": True,
        "role_count_min": 1,
        "role_count_max": 3,
        "role_usage_mode": "supporting-only",
        "role_density_preference": role_scale,
        "role_scene_relationship": choose_role_scene_relationship(scene_concept, role_scale, index),
        "role_in_scene_mode": role_in_scene_mode,
        "role_in_scene_guidance": choose_role_in_scene_guidance(role_in_scene_mode, page_type, scene_concept),
        "reserved_zone": RESERVED_ZONE,
        "reserved_zone_rules": RESERVED_ZONE_RULES,
        "pagination_mode": page.get("pagination_mode", "explicit"),
        "section_kind": page.get("section_kind", "manual"),
        "step_label": page.get("step_label", ""),
        "step_index": page.get("step_index"),
        "step_page_index": page.get("step_page_index", 1),
        "step_page_total": page.get("step_page_total", 1),
        "is_continued": page.get("is_continued", False),
        "forbid_visible_page_number": True,
        "role_action": role_action,
        "role_expression": role_expression,
        "role_identity_lock": "strict",
        "role_scale": role_scale,
        "role_position": role_position,
        "role_framing": role_framing,
        "role_composition_mode": role_composition_mode,
        "camera_energy": camera_energy,
        "role_pose_hint": role_pose_hint,
        "role_variation_guard": page.get("role_variation_guard", "semantic-match-primary"),
        "role_variation_reason": page.get("role_variation_reason", "基础语义分配，后续允许按整套选题做去重调整"),
        "icon_set": icon_set,
        "layout_constraints": DEFAULT_LAYOUT_CONSTRAINTS,
    }


def build_prompt(spec: dict[str, Any], role: dict[str, Any]) -> str:
    line_color_desc = "black lines" if spec["line_color"] == "black" else "white lines"
    title = spec.get("display_title") or spec.get("title_text") or spec["page_title"]
    source_title_full = spec.get("source_title_full") or spec.get("title_text") or spec["page_title"]
    hook = spec.get("page_hook", "")
    cue_phrases = spec.get("cue_phrases") or spec.get("talking_points") or []
    flow_labels = spec.get("flow_labels") or []
    support_labels = spec.get("support_labels") or []
    subtitle_notes = spec.get("subtitle_notes") or []
    primary_info_blocks = spec.get("primary_info_blocks") or []
    secondary_info_blocks = spec.get("secondary_info_blocks") or []
    icons = " / ".join(spec["icon_set"])
    visual_keywords = " / ".join(spec.get("visual_keywords") or [])
    visual_action_cues = " / ".join(spec.get("visual_action_cues") or [])
    visual_scene_cues = " / ".join(spec.get("visual_scene_cues") or [])
    role_action_tags = " / ".join(spec.get("role_action_tags") or [])
    role_expression_tags = " / ".join(spec.get("role_expression_tags") or [])
    micro_visuals = " / ".join(spec.get("supporting_micro_visuals") or [])
    role_prompt = role.get("prompt_block") or role["meta"].get("display_name", role["slug"])
    scale_ratio = role_scale_ratio(spec.get("role_scale", "medium-support"))
    density = spec.get("density_profile") or {}
    cue_phrase_limit = density.get("talking_points", 3)
    micro_visual_limit = density.get("micro_visuals", 4)
    support_label_limit = density.get("support_labels", 2)
    continued_hint = ""
    if spec.get("is_continued"):
        continued_hint = (
            "- this page continues the same step internally, but do not render any page number, step sequence, x/y counter, badge, or continuation marker on the slide\n"
        )
    cue_phrases_block = "\n".join(f"  - {point}" for point in cue_phrases) if cue_phrases else "  - 无"
    flow_labels_block = "\n".join(f"  - {point}" for point in flow_labels) if flow_labels else "  - 无"
    support_labels_block = "\n".join(f"  - {point}" for point in support_labels) if support_labels else "  - 无"
    subtitle_notes_block = "\n".join(f"  - {point}" for point in subtitle_notes) if subtitle_notes else "  - 无"
    primary_blocks_block = "\n".join(f"  - {point}" for point in primary_info_blocks) if primary_info_blocks else "  - 无"
    secondary_blocks_block = "\n".join(f"  - {point}" for point in secondary_info_blocks) if secondary_info_blocks else "  - 无"
    face_anchor_ref, action_ref, action_reason = split_role_reference_chain(spec, role, preferred_variant="standard")
    face_anchor_label = face_anchor_ref.get("name", "") if face_anchor_ref else "none"
    action_ref_label = action_ref.get("name", "") if action_ref else "none"
    return f"""---
aspect_ratio: "16:9"
---
Generate one standalone 16:9 horizontal Chinese PPT page as a full-slide visual note page.

Style:
- hand-drawn visual-note style
- neutral, concise, clear
- monochrome line-art look with restrained but layered handwritten Chinese text
- exact background color: {spec["background_color"]}
- use {line_color_desc} for the main drawing and handwritten annotations
- visible red accents are REQUIRED on every page
- use exactly {spec.get("red_accent_count_min")} to {spec.get("red_accent_count_max")} clearly visible red accent areas
- red accents may only appear on: {" / ".join(spec.get("red_accent_allowed_targets") or [])}
- use red accents to mark key words, key arrows, tags, result blocks, and warning points
- do not use red as a large background area or the main page color
- no dense layout, no realistic rendering, no glossy 3D
- do not simplify this into a minimal poster or ultra-clean infographic
- keep the page readable, visual, and graph-first; text should be light and source-anchored

Scene direction:
- scene priority: {spec.get("scene_priority")}
- scene concept: {spec.get("scene_concept")}
- scene layout type: {spec.get("scene_layout_type")}
- this page should feel like one standalone poster-like page
- the title is part of the composition, not a fixed template block
- title layout mode: {spec.get("title_layout_mode")}
- title alignment: {spec.get("title_alignment")}
- title anchor zone: {spec.get("title_anchor_zone")}
- title flow direction: {spec.get("title_flow_direction")}
- title layout guidance: {spec.get("title_layout_guidance")}
- title placement should follow the page scene
- do not force the title into the same position across pages
- role scene relationship: {spec.get("role_scene_relationship")}
- role in scene mode: {spec.get("role_in_scene_mode")}
- role in scene guidance: {spec.get("role_in_scene_guidance")}
- supporting micro visuals: {micro_visuals}
- source-text-only visual mode: {spec.get("source_text_only_visual_mode")}
- inferred visuals forbidden: {spec.get("inferred_visuals_forbidden")}
- do not invent objects, results, metaphors, or industry scenes that do not appear in the source text
- if the source text does not mention a scene, keep the page abstract and structural instead of adding one
- visible text may only come from the page title, source-rooted cue phrases, flow labels, support labels, and short subtitle notes
- do not add explanatory captions, summary lines, slogans, or commentary that go beyond the source wording
- prioritize scene, structure, diagram, relation map, process map, contrast map, and object interaction over text blocks
- if an idea can be understood from the drawing, arrows, structure, or scene, do not restate it as extra text
- every page must have one dominant graphic structure and at least 2 supporting visual elements
- role count range: {spec.get("role_count_min")} to {spec.get("role_count_max")}
- role usage mode: {spec.get("role_usage_mode")}
- role density preference: {spec.get("role_density_preference")}
- do not redesign a new character from scratch
- face anchor only: {face_anchor_label}
- action reference only: {action_ref_label}
- action reference reason: {action_reason or "none"}
- use the face anchor only to lock face shape, hairline, glasses outline, age feel, and face proportions
- use the action reference only to borrow action, gesture, expression tendency, or pose direction
- other references must not redefine the face
- do not use anchor photos, pose images, external character photos, or API-side identity references
- reserved area spec: {spec.get("reserved_zone")}
- reserved area rules: {", ".join(spec.get("reserved_zone_rules") or [])}
- keep the bottom-left 13cm by 13cm area as uninterrupted pure background color
- do not draw any frame, border, placeholder box, guide shape, label, or scene edge around that area
- build the full page around one unified main scene driven by the page title and theme
- do not start from a generic host figure and then decorate around him
- the character should be embedded into the scene, not a separate hero poster
- the page may also rely on non-character illustrations, structure sketches, diagrams, arrows, and relation maps; the character is only one element, not the whole page
- do not place the character outside the main scene as a separate host
- do not place the character in an empty side area while the main content sits elsewhere
- do not let the character act as a floating presenter detached from the page structure
- the character must physically interact with the main scene, labels, props, arrows, tools, or structure
- if the character cannot interact with the main structure at a large size, reduce the character size and move him into the structure instead of isolating him
- every page should read as one integrated scene, not character plus decorations
- the character is not required to be the visual center; the scene, title, structure, and theme can carry the page
- if the main scene is already clear, keep the character smaller and less central
- when the page mainly expresses a process, object, or relationship, build that structure first and embed 1 to 3 supporting characters afterward
- never use 4 or more same-role characters on one page
- do not repeat the same role more than {spec.get("role_count_max")} times on one page
- the last page should favor a richer structure graphic over another large character portrait

Recurring character:
{role["meta"].get("display_name", role["slug"])} appears on this page.
{role_prompt}

Character action:
{spec["role_action"]}

Character expression:
{spec["role_expression"]}

Character staging:
- role scale: {spec.get("role_scale")}
- character occupies {scale_ratio}
- role position: {spec.get("role_position")}
- role framing: {spec.get("role_framing")}
- composition mode: {spec.get("role_composition_mode")}
- camera energy: {spec.get("camera_energy")}
- pose hint: {spec.get("role_pose_hint")}
- expression tags: {role_expression_tags}
- same-topic variation guard: {spec.get("role_variation_guard")}
- same-topic variation reason: {spec.get("role_variation_reason")}
- role count min/max: {spec.get("role_count_min")} / {spec.get("role_count_max")}
- role usage mode: {spec.get("role_usage_mode")}
- this page may include 1 to 3 instances of the same role, but never a crowd of avatars

Character constraints:
- this is not a new character design task
- only change action, hand gesture, facial expression intensity, and body pose
- do not change face identity, hairstyle, glasses, outfit, age feel, or drawing style
- keep the same slicked-back side-part black hair, black rectangular glasses, white blazer, and black inner shirt on every page
- the hair must stay smooth, flatter, tighter to the scalp, and cleaner like the top-left face-anchor look
- skin tone base should stay close to {spec.get("skin_tone_base")}
- allow only very small same-family shading shifts around that base for normal hand-drawn volume
- do not globally lighten, pinken, gray out, cool down, or switch to a pale-skin template
- if the skin drifts away from the #e8a668 family toward a whiter commercial illustration skin, treat it as a failed render and redo
- when showing full body, use jeans and clean white sneakers
- keep the face shape oval and slightly fuller, not slim or narrow
        - for frontal or near-frontal views, keep the eye structure consistent with the approved face-anchor identity
- avoid reducing frontal eyes to tiny dot eyes
- keep clearer eye shape and gaze inside the glasses when the face is front-facing
- side views or highly exaggerated expressions may simplify the eyes slightly, but frontal identity should stay stronger
- push the facial expression and body language to a clearly exaggerated presentation level
- use more expressive eyebrows, clearer mouth shapes, and wider open-hand gestures
- the character must still be clearly recognizable as the same person
- chibi or more cute-proportion rendering is allowed, but identity drift is not allowed
- do not replace the blazer with a hoodie, sweatshirt, coat, T-shirt outer layer, or casualwear silhouette
- do not change the hairstyle away from the slicked-back side-part reference; no front-combed variant, no curly hair, no fluffy hair
- do not add wavy texture, curled strands, airy volume, or a visibly raised pompadour-like top
- if the hair is back-combed but still looks curly, wavy, fluffy, airy, or too tall, treat it as failed identity lock and redo
        - do not change the glasses shape, face shape, jawline, hairline, age feel, or skin-tone identity locked by the face anchor
        - when the model tries to make the character cuter, still prioritize the face anchor identity over free cartoon redesign or anchor-photo reinterpretation
        - text instructions may refine acting and composition only; they must not redefine face identity away from the face anchor
- do not force the microphone on every page
- if a microphone appears, the mic flag may show the Chinese text "知富星球"
- do not keep the character at the same size across consecutive pages
- do not keep the character centered on every page
- page 1 is not automatically a cover; only an explicit cover page may use a more centered hero composition
- later pages should actively vary left, center, right placement and viewing distance
- match the character action to the page title scene instead of using a generic hosting pose
- within the same topic, do not repeat a near-identical presenter pose, facial expression, or action family from adjacent pages
- if a neighboring page already uses a similar action idea, switch this page to a different acting pattern, expression family, or staging choice
- the character should be part of the page mechanism, not a separate commentator
- the character must touch, use, enter, cross, operate, compare, or be constrained by the main scene
- use only the face anchor to lock the approved face shape, hair direction, glasses structure, and skin tone
- use the selected action reference only to borrow gesture, action direction, or pose energy
- do not use any anchor photo, pose library image, external character photo, or API-side identity reference
- decide action from the page text and composition fields while preserving the face-anchor character identity
- the character must stay out of the bottom-left pure-background no-draw area
- that bottom-left area must remain plain background only, with no edges, marks, scene fragments, labels, or partial props
- that bottom-left area is not a "best effort" suggestion; it is a hard no-draw area
- use the visual action and scene cues to decide what the character is doing
- visual action cues and scene cues must stay source-anchored and may only restate explicit script content, not invent new meaning
- visual action cues and scene cues are for image choreography only, not for extending the script
- the page may use a small number of short support labels only when they are directly rooted in the source text
- support labels must stay short, restrained, and structural; they must not turn into paragraphs or a label wall
- keep text-light expression by default: most pages should rely on graphics first and use only a very small number of short handwritten text anchors
- do not create extra note-like statements once the scene, structure, and graphic relationship already express the point
- every independent small visual unit, icon group, mistake card, process node, relation block, or comparison unit must have one nearby short handwritten label
- each small visual label should stay within roughly 2 to 6 Chinese characters and should preferably come from the source wording
- if the page contains several important mini-visuals but some of them have no text anchor, treat the page as incomplete and add short source-rooted labels

Page content:
- display title mode: {spec.get("display_title_mode")}
- page title shown on the page should be: {title}
- full source title for reference: {source_title_full}
- title strategy: {spec.get("title_strategy")}
- if title strategy is explicit-step-with-subtitle, preserve the first step layer as the main title and preserve subtitle notes as the second step layer
- if title strategy is explicit-step-with-subtitle, cue phrases and support labels must not replace or override the main title or subtitle
- if title strategy is extracted-full-hook, the title must stay a complete standalone judgment sentence, not a fragment, not a noun phrase, and not half a sentence
- deck title: {spec.get("deck_title") or ""}
- this page is for speaking or recording support, not for displaying the full script
- do NOT render the full transcript or a large paragraph box
- show at most {cue_phrase_limit} short handwritten Chinese cue phrases
- show at most {support_label_limit} short structural support labels
- use only 1 to {micro_visual_limit} supporting micro visuals or callout elements
- use a graph-first, text-light, source-anchored page mode
- visible cue phrases should stay close to the source script
- support labels may only be short source-rooted structure words; do not rewrite them into new explanatory language
- support labels must be extremely short, must serve the diagram, and must not become long explanations, slogans, or filler text
- if a small cue phrase cannot be rendered accurately, omit it instead of fabricating text
- use sparse guiding words instead of full sentences
- never let text become the main explanation layer of the page
- every major mini-visual should pair with one short label so the speaker can point to it while explaining
- do not render any page number, slide number, x/y counter, chapter counter, corner badge, or sequence marker such as 4/16 or 8/16 anywhere on the image
- page hook: {hook}
- visible cue phrases:
{cue_phrases_block}
- visible flow labels:
{flow_labels_block}
- allowed support labels:
{support_labels_block}
- primary info blocks:
{primary_blocks_block}
- primary info blocks are only minimal anchor points; the page meaning must mainly come from the visual structure
- do not expand primary blocks into text-heavy rectangles
- secondary info blocks:
{secondary_blocks_block}
- secondary info blocks should normally stay empty; only use them when absolutely necessary for a complex page
- subtitle notes for long-title pages:
{subtitle_notes_block}
- if subtitle notes are present for a step-structured page, render them as the secondary title layer under or near the main title, not as loose body labels
- text density mode: {spec.get("text_density_mode")}
- page type: {spec["page_type"]}
- section kind: {spec.get("section_kind", "manual")}
- supporting icons for drawing only: {icons}
- visual keywords for scene planning only, derived from the source text only: {visual_keywords}
- visual action cues for image only, derived from the source text only: {visual_action_cues}
- visual scene cues for image only, derived from the source text only: {visual_scene_cues}
- role action tags: {role_action_tags}
- role expression tags: {role_expression_tags}
{continued_hint}

Layout constraints:
- make the title large and prominent
- supporting words should be fewer, smaller, and subordinate to the drawing
- build the page around the main scene, diagram, or relation structure first, then add only the minimum source-rooted text anchors needed for fast understanding
- leave the bottom-left corner free of any text
- leave the bottom-left corner free of any characters, icons, lines, arrows, or scene borders
- keep the slide readable on a 16:9 PPT page
- this is one full-slide image, not a collage of many cards
- no large paragraph box anywhere on the page
- no label wall, no annotation cloud, no fake handwritten commentary block
- do not collapse the page into only a title plus one or two words with excessive empty space unless the source text itself is extremely sparse
- when the page is structure-first, let the character become a supporting accent instead of the visual center
- when the page is emotion-first, let the character become the stronger visual subject
- if the main scene is already clear, remove extra callouts instead of adding more
- keep the total amount of visible text blocks and micro callouts restrained
- prefer replacing extra words with icons, object relations, arrows, process steps, contrast shapes, and scene metaphors that stay source-anchored
- keep the bottom-left 13cm by 13cm area as natural continuous background color only
- inside that area: no text, no icons, no character, no lines, no arrows, no footprints, no pit edges, no decorative marks
- do not draw any box, grid cell, frame, boundary, or placeholder for that area
- compose the main scene around this no-draw area from the start

Text rendering:
- all visible text should look handwritten
- for exact-title pages, preserve the title exactly as provided
- for compressed-title pages, keep any leading guided number such as 第一点 / 第1步 / 第1条 / 认知一 inside the main title, and let subtitle notes carry only the remaining meaning
- do not render large paragraph text, label walls, fake commentary, or long explanatory sentences
- visible text should stay limited to the display title, cue phrases, flow labels, subtitle notes, and short support labels above
- the page should still feel understandable if most viewers first read the picture and only then notice the text
- this prompt is intended for Codex image generation only; external API or compatible-model rendering is forbidden
"""


def build_ppt_from_pages(page_paths: list[pathlib.Path], out_path: pathlib.Path) -> None:
    prs = Presentation()
    prs.slide_width = 16256000
    prs.slide_height = 9144000
    blank = prs.slide_layouts[6]
    for page_path in page_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(page_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    ensure_dir(out_path.parent)
    prs.save(str(out_path))


def load_specs_from_workdir(workdir: pathlib.Path) -> list[dict[str, Any]]:
    specs_path = workdir / "pages-spec.json"
    if not specs_path.exists():
        raise FileNotFoundError(f"找不到 pages-spec.json: {specs_path}")
    payload = json.loads(specs_path.read_text(encoding="utf-8"))
    if isinstance(payload, list):
        return payload
    if isinstance(payload, dict) and isinstance(payload.get("pages"), list):
        return payload["pages"]
    raise ValueError(f"pages-spec.json 结构不合法: {specs_path}")


def build_cover_prompt(cover_spec: dict[str, Any], role: dict[str, Any]) -> str:
    support_labels = cover_spec.get("support_labels") or []
    cue_phrases = cover_spec.get("cue_phrases") or []
    support_labels_block = "\n".join(f"  - {item}" for item in support_labels) if support_labels else "  - 无"
    cue_phrases_block = "\n".join(f"  - {item}" for item in cue_phrases) if cue_phrases else "  - 无"
    face_anchor_image = cover_spec.get("role_face_anchor_image") or ""
    action_reference_image = cover_spec.get("role_action_reference_image") or ""
    action_reference_reason = cover_spec.get("role_action_reference_reason") or ""
    reference_block = action_reference_reason or "ip-face-anchor-top-left"
    role_prompt = role.get("prompt_block") or role["meta"].get("display_name", role["slug"])
    return f"""---
aspect_ratio: "{cover_spec['aspect_ratio']}"
---
Generate one standalone Chinese cover image in {cover_spec['aspect_ratio']} ratio.

Style:
- same hand-drawn visual-note style as the existing page images
- exact background color: {cover_spec['background_color']}
- use black lines for the drawing and handwritten annotations
- visible red accents are REQUIRED on the cover
- use 2 to 4 clearly visible red accent areas on keywords, arrows, tags, warning points, or result blocks only
- keep the same role identity, same page-world feeling, and same rough notebook-energy as the page set
- do not drift into glossy poster art, realistic portrait art, or generic internet cover style

Cover goal:
- this is an extra cover asset, not page 1 of the main slide set
- the visible main title must be exactly: {cover_spec['title_text']}
- the character must appear and remain one of the core visual subjects
- the character action and expression must clearly match the title meaning
- use a same-theme adaptation approach: keep the same topic expression, but compose this ratio specifically instead of simply cropping another layout
- composition guidance: {cover_spec['composition_guidance']}
- title zone guidance: {cover_spec['title_zone']}
- cover decorative modules must stay between {cover_spec.get("cover_accent_blocks_min")} and {cover_spec.get("cover_accent_blocks_max")} blocks total
- count cue phrases and support labels together as one decorative-block pool
- decorative modules may assist the title, but they must not become a label wall, annotation cloud, or dense ring of notes
- role count range is {cover_spec.get("role_count_min")} to {cover_spec.get("role_count_max")}
- role usage mode is {cover_spec.get("role_usage_mode")}
- source-text-only visual mode is mandatory for the cover
- do not invent objects, results, metaphors, or industry scenes that are not explicit in the title or source phrases
- the cover may give the character slightly more presence than a正文页, but the character still counts as a supporting element and must stay within 1 to 3 instances
- do not use repeated avatar crowds or a row of many same-role figures to fill the page
- do not redesign a new character from scratch
- face anchor image: {face_anchor_image or "none"}
- action reference image: {action_reference_image or "none"}
- action reference reason: {reference_block}
- use the face anchor only to lock the person; use the action reference only to borrow action or expression tendency

Recurring character:
{role["meta"].get("display_name", role["slug"])} appears on this cover.
{role_prompt}

Character constraints:
        - use exactly the same role chain as the page images: same role slug, same face-anchor identity, same reference-image path family, and same action-reference library
- keep the same slicked-back side-part black hair, black rectangular glasses, white blazer, and black inner shirt
- the hair must stay smooth, flatter, tighter to the scalp, and cleaner like the top-left face-anchor look
- keep the same face identity and warmer, slightly deeper skin tone already locked by the approved role references
- skin tone base should stay close to {cover_spec.get("skin_tone_base")}
- allow only very small same-family shading shifts around that base for normal hand-drawn volume
- do not globally lighten, pinken, gray out, cool down, or switch to a pale-skin template
- if the skin drifts away from the #e8a668 family toward a whiter commercial illustration skin, treat it as a failed render and redo
- keep the cover character expressive, dynamic, and clearly recognizable as the same person
- do not change the outfit, hair direction, glasses shape, or face identity
- keep the face shape oval and slightly fuller, not slim or narrow
        - for frontal or near-frontal views, keep the eye structure consistent with the approved face-anchor identity
- avoid reducing frontal eyes to tiny dot eyes
- keep clearer eye shape and gaze inside the glasses when the face is front-facing
- preferred character reference poses: {reference_block}
- only face-anchor reference image: {face_anchor_image or "none"}
- only action-reference image: {action_reference_image or "none"}
- character reference role: {reference_block}
        - do not loosen the face lock just because this is a cover
- do not switch to another face, another hairstyle family, another glasses shape, another clothing identity, or another reference path
- do not add wavy texture, curled strands, airy volume, or a visibly raised pompadour-like top
        - do not invent a separate cover-specific character variant
        - all facial identity, hairstyle, glasses, skin tone, outfit, and body-language constraints must match the page-image character rules exactly
        - do not change the glasses shape, face shape, jawline, hairline, age feel, or skin-tone identity locked by the nine-grid calibration board
        - when the model tries to make the character cuter, still prioritize the face-anchor identity over free cartoon redesign
        - because the cover usually uses a larger character, prioritize the face-anchor identity even more strictly than ordinary pages
        - if forced to choose, preserve the face-anchor identity over extra stylization
        - face-anchor image: {face_anchor_image or "none"}
        - action-reference role: {reference_block}
        - use the face anchor as the only identity source for the approved face shape, hair direction, glasses structure, and skin tone
        - do not use any anchor photo, pose library image, external character photo, or API-side identity reference
        - decide action from the page text and composition fields while preserving the face-anchor character identity
        - text instructions may refine acting and composition only; they must not redefine face identity away from the calibration board

Character acting:
- action guidance: {cover_spec['action_guidance']}
- expression guidance: {cover_spec['expression_guidance']}
- role action tags: {" / ".join(cover_spec.get("role_action_tags") or [])}
- role expression tags: {" / ".join(cover_spec.get("role_expression_tags") or [])}

Visible text:
- exact main title: {cover_spec['title_text']}
- optional short cue phrases:
{cue_phrases_block}
- optional short support labels:
{support_labels_block}
- total decorative modules on the cover must not exceed {cover_spec.get("cover_accent_blocks_max")} blocks
- support labels must stay very short and structural
- support labels may include a small amount of short English words for visual clarity
- do not turn the support labels into long English sentences or a mixed-language text wall
- do not render paragraphs, fake commentary blocks, or dense text
- do not render any page number, slide number, x/y counter, or corner badge

Composition:
- prioritize one strong, readable cover composition
- make the title large and readable
- use the character pose, gaze, and gesture to reinforce the title conflict
- 3:4 should feel more vertical and centered; 4:3 should feel wider and more lateral
- a few tiny accent marks, arrows, coins, tags, or symbols are allowed when they reinforce the topic
- keep the final image bold, clean, and high-contrast on the yellow base

Text rendering:
- all visible text should look handwritten
- preserve the exact title as written
- if small labels cannot be rendered accurately, omit them instead of inventing text
- this prompt is intended for Codex image generation only; external API or compatible-model rendering is forbidden
"""


def ensure_page_images_exist(page_paths: list[pathlib.Path]) -> None:
    missing = [str(path) for path in page_paths if not path.exists()]
    if missing:
        raise FileNotFoundError("以下目标归档图片未写回，任务未完成，暂时不能装配 PPT:\n" + "\n".join(missing))


def ensure_archive_writeback_complete(page_paths: list[pathlib.Path], archive_dir: pathlib.Path) -> None:
    missing = [str(path) for path in page_paths if not path.exists()]
    outside_archive = [str(path) for path in page_paths if archive_dir not in path.parents]
    if missing:
        raise FileNotFoundError("以下目标归档图片未写回，任务未完成:\n" + "\n".join(missing))
    if outside_archive:
        raise FileNotFoundError("以下图片路径不在指定归档目录内，不能视为正式交付结果:\n" + "\n".join(outside_archive))


def build_codex_handoff(
    specs: list[dict[str, Any]],
    cover_outputs: list[dict[str, Any]],
    selected_pages: list[int],
    workdir: pathlib.Path,
    prompt_paths: list[pathlib.Path],
    page_paths: list[pathlib.Path],
    cover_prompt_paths: dict[str, pathlib.Path],
    role: dict[str, Any],
) -> tuple[pathlib.Path, pathlib.Path, pathlib.Path, pathlib.Path]:
    prompts_by_page = {spec["page_index"]: prompt_path for spec, prompt_path in zip(specs, prompt_paths)}
    pages_by_page = {spec["page_index"]: page_path for spec, page_path in zip(specs, page_paths)}
    selected_specs = [spec for spec in specs if spec["page_index"] in selected_pages]
    items = []
    for spec in selected_specs:
        face_anchor_ref, action_ref, action_reason = split_role_reference_chain(spec, role, preferred_variant="standard")
        item = {
            "page_index": spec["page_index"],
            "page_title": spec["page_title"],
            "title_text": spec.get("title_text"),
            "display_title_mode": spec.get("display_title_mode"),
            "display_title": spec.get("display_title"),
            "source_title_full": spec.get("source_title_full"),
            "deck_title": spec.get("deck_title"),
            "page_hook": spec.get("page_hook"),
            "cue_phrases": spec.get("cue_phrases"),
            "flow_labels": spec.get("flow_labels"),
            "support_labels": spec.get("support_labels"),
            "subtitle_notes": spec.get("subtitle_notes"),
            "talking_points": spec.get("talking_points"),
            "text_density_mode": spec.get("text_density_mode"),
            "visual_keywords": spec.get("visual_keywords"),
            "visual_action_cues": spec.get("visual_action_cues"),
            "visual_scene_cues": spec.get("visual_scene_cues"),
            "role_action_tags": spec.get("role_action_tags"),
            "role_expression_tags": spec.get("role_expression_tags"),
            "role_scale": spec.get("role_scale"),
            "role_position": spec.get("role_position"),
            "role_framing": spec.get("role_framing"),
            "role_composition_mode": spec.get("role_composition_mode"),
            "role_in_scene_mode": spec.get("role_in_scene_mode"),
            "role_in_scene_guidance": spec.get("role_in_scene_guidance"),
            "camera_energy": spec.get("camera_energy"),
            "role_pose_hint": spec.get("role_pose_hint"),
            "role_identity_lock": spec.get("role_identity_lock"),
            "role_variation_guard": spec.get("role_variation_guard"),
            "role_variation_reason": spec.get("role_variation_reason"),
            "section_kind": spec.get("section_kind"),
            "step_label": spec.get("step_label"),
            "step_page_index": spec.get("step_page_index"),
            "step_page_total": spec.get("step_page_total"),
            "forbid_visible_page_number": spec.get("forbid_visible_page_number"),
            "background_color": spec["background_color"],
            "line_color": spec["line_color"],
            "scene_priority": spec.get("scene_priority"),
            "scene_concept": spec.get("scene_concept"),
            "scene_layout_type": spec.get("scene_layout_type"),
            "title_layout_mode": spec.get("title_layout_mode"),
            "title_alignment": spec.get("title_alignment"),
            "title_anchor_zone": spec.get("title_anchor_zone"),
            "title_flow_direction": spec.get("title_flow_direction"),
            "supporting_micro_visuals": spec.get("supporting_micro_visuals"),
            "role_scene_relationship": spec.get("role_scene_relationship"),
            "role_required": spec.get("role_required"),
            "source_text_only_visual_mode": spec.get("source_text_only_visual_mode"),
            "inferred_visuals_forbidden": spec.get("inferred_visuals_forbidden"),
            "role_count_min": spec.get("role_count_min"),
            "role_count_max": spec.get("role_count_max"),
            "role_usage_mode": spec.get("role_usage_mode"),
            "role_density_preference": spec.get("role_density_preference"),
            "reserved_zone": spec.get("reserved_zone"),
            "reserved_zone_rules": spec.get("reserved_zone_rules"),
            "role_face_anchor_image": str(face_anchor_ref["file_path"]) if face_anchor_ref else "",
            "role_action_reference_image": str(action_ref["file_path"]) if action_ref else "",
            "role_action_reference_reason": action_reason or (action_ref.get("name", "") if action_ref else ""),
            "role_reference_strategy": "anchor-plus-pose",
            "prompt_file": str(prompts_by_page[spec["page_index"]]),
            "output_image": str(pages_by_page[spec["page_index"]]),
            "output_image_required": True,
            "archive_writeback_required": True,
            "delivery_mode": "direct-or-copy-then-clean",
        }
        item.update(build_reference_lock_metadata())
        items.append(item)
    is_full_run = len(selected_pages) == len(specs)
    cover_items = []
    if is_full_run:
        cover_ref_source = selected_specs[0] if selected_specs else (specs[0] if specs else {})
        face_anchor_ref, action_ref, action_reason = build_cover_reference_chain(cover_ref_source, role)
        for cover_spec in cover_outputs:
            cover_item = {
                "cover_type": cover_spec["cover_type"],
                "aspect_ratio": cover_spec["aspect_ratio"],
                "title_text": cover_spec["title_text"],
                "display_title": cover_spec["display_title"],
                "background_color": cover_spec["background_color"],
                "line_color": cover_spec["line_color"],
                "style_mode": cover_spec["style_mode"],
                "role_required": cover_spec["role_required"],
                "support_labels_allowed": cover_spec["support_labels_allowed"],
                "cue_phrases": cover_spec.get("cue_phrases"),
                "support_labels": cover_spec.get("support_labels"),
                "cover_accent_blocks_min": cover_spec.get("cover_accent_blocks_min"),
                "cover_accent_blocks_max": cover_spec.get("cover_accent_blocks_max"),
                "cover_accent_blocks_total": cover_spec.get("cover_accent_blocks_total"),
                "role_action_tags": cover_spec.get("role_action_tags"),
                "role_expression_tags": cover_spec.get("role_expression_tags"),
                "expression_guidance": cover_spec.get("expression_guidance"),
                "action_guidance": cover_spec.get("action_guidance"),
                "composition_guidance": cover_spec.get("composition_guidance"),
                "title_zone": cover_spec.get("title_zone"),
                "source_text_only_visual_mode": cover_spec.get("source_text_only_visual_mode"),
                "inferred_visuals_forbidden": cover_spec.get("inferred_visuals_forbidden"),
                "role_count_min": cover_spec.get("role_count_min"),
                "role_count_max": cover_spec.get("role_count_max"),
                "role_usage_mode": cover_spec.get("role_usage_mode"),
                "role_density_preference": cover_spec.get("role_density_preference"),
                "role_face_anchor_image": str(face_anchor_ref["file_path"]) if face_anchor_ref else "",
                "role_action_reference_image": str(action_ref["file_path"]) if action_ref else "",
                "role_action_reference_reason": action_reason or (action_ref.get("name", "") if action_ref else ""),
                "role_reference_strategy": "anchor-plus-pose",
                "prompt_file": str(cover_prompt_paths[cover_spec["cover_type"]]),
                "output_image": cover_spec["output_image"],
                "output_image_required": True,
                "archive_writeback_required": True,
                "delivery_mode": "direct-or-copy-then-clean",
            }
            cover_item.update(build_reference_lock_metadata())
            cover_items.append(cover_item)
    handoff = {
        "mode": "codex",
        "workdir": str(workdir),
        "work_package_dir": str(workdir),
        "archive_dir": str(page_paths[0].parent if page_paths else ""),
        "final_archive_dir": str(page_paths[0].parent if page_paths else ""),
        "output_image_required": True,
        "archive_writeback_required": True,
        "delivery_mode": "direct-or-copy-then-clean",
        "selected_pages": selected_pages,
        "full_run": is_full_run,
        "pages": items,
        "covers": cover_items,
    }
    handoff.update(build_reference_lock_metadata())
    handoff_path = workdir / "codex-handoff.json"
    workflow_path = workdir / "codex-workflow.txt"
    write_json(handoff_path, handoff)
    workflow = (
        "Codex 出图工作流\n"
        "阶段 1：参考图接入校验\n"
        "0. 正式出图前先确认 config/skill-config.json 里的 archive_root 已显式配置；默认工作包固定写入 `outputs/<选题名>-work`，正式成品固定写入 `archive_root/<选题名>/`，若 archive_root 为空且会触发 skill 相邻目录回退，必须先停止并修正配置，再继续正式出图。\n"
        "1. 出图前先核对 codex-handoff.json 顶层的 reference_lock_required、text_only_render_forbidden、renderer_must_attach_reference_images、identity_reference_mode、identity_baseline_image、fixed_identity_traits_required、secondary_identity_sources_forbidden、required_reference_inputs、optional_non_identity_references、fail_if_reference_images_not_actually_attached、codex_only_rendering_required、external_api_rendering_forbidden、allowed_render_mode。\n"
        "2. 每次新文案正式出图前，都必须重新让模型真实看到 role_face_anchor_image 这张主脸锚点图；不能只复用上一次会话记忆。\n"
        "3. 固定人物特征文字只作为辅助稳定器，帮助锁定黑框眼镜、暖肤色、背头侧分黑发、偏长椭圆脸、白西装黑内搭、全身时蓝色牛仔裤和白鞋、手绘卡通线稿风；这些文字不能替代主脸锚点图，也不能变成第二身份源。\n"
        "4. 左上正脸主脸模板负责锁定额头高度、发际线方向、五官比例、眼镜大小、脸宽脸长比和年龄感；动作参考单图只允许提供角度和动作变化边界，不得重新定义主脸。\n"
        "5. 本流程禁止真人锚点照、姿态库身份图、第二角色参考图进入正式身份链路；人物身份只审核 role_face_anchor_image 这张主脸锚点图。\n"
        "6. 只看到参考图路径、只把参考要求写进文字 prompt、只抽取人物特征文字，或无法确认主脸锚点图是否真的被模型使用，全部视为未通过校验。\n"
        "7. 若当前工具只支持文字 prompt、不支持真实参考图输入，必须立即停止并汇报：当前执行环境无法真实挂载角色主脸锚点参考图，本次停止生成，未进入正式出图链路。\n"
        "8. 不允许先出图再人工挑一张像的，也不允许改用泛化人物继续生成。\n"
        "9. 同一页若出现多个同角色小人、不同景别或不同动作，也必须继续使用同一张主脸锚点图锁定；若当前执行端做不到稳定锁定，必须先降为单主角色或单次校准，不允许硬出整页。\n"
        "阶段 2：逐页出图与归档\n"
        "10. 只有通过阶段 1 校验后，才按 codex-handoff.json 中的 page_index 顺序逐页生成图片；默认先在工作包目录生成清单，再把成品直接写回正式归档目录。\n"
        "11. 每生成一页，必须先回头对照主脸锚点图确认是不是同一个人；先做左上正脸主模板一致性判定，再核对眼镜、肤色、发型、脸型、年龄感和白西装黑内搭。不一致就只重做该页。\n"
        "12. 每生成一页，优先尝试把最终文件直接写入该页 handoff 里的 output_image；如果当前平台只能先产出中间 PNG，再复制归档，也允许继续。\n"
        "13. 若走复制归档链路，必须先把中间 PNG 复制到正式目录对应文件名，再删除这张已归档成功的缓存图；中间目录不得作为长期交付目录保留。\n"
        "13.1 默认按单文件粒度清理缓存：只删除本轮刚刚归档成功的那一张缓存图，不清空整包缓存目录，也不删除其他历史缓存文件。\n"
        "14. 若当前平台既不能直接写正式目录，也不能先得到可复制的中间 PNG，才视为真正阻塞并立即停止。\n"
        "15. 仅存在于临时生成目录、缓存目录、对话附件或 C 盘缓存中的图片都不算完成；工作包目录和中间图片目录都不是最终交付目录。\n"
        "16. 如果这次是整套完整出图，在正文页全部归档到正式目录后，还要继续生成 cover-3x4.png 和 cover-4x3.png。\n"
        "17. 禁止脱离 handoff 另起一套临时手写 prompt 直接出图。\n"
        "18. 先做 1-2 页校准风格，再继续后续页面；试投页必须先通过左上正脸主模板一致性，再看构图。\n"
        "19. 页面可见文字只认 handoff 里的 title_text、display_title、cue_phrases、flow_labels、subtitle_notes 和 source-text-only 结构词，不得补成解释型长句。\n"
        "20. 左下角 13cm × 13cm 只能保留自然连续的纯背景，不画框、不画格、不画占位提示。\n"
        "21. 页面允许同时使用非人物插画、结构示意图、流程图、箭头和关系图，但这些元素必须直接回指到原文，不得额外补出金币、火箭、收入回报、办公室工作等原文未说的东西。\n"
        "22. 人物必须先嵌入主结构再谈造型，不能作为画外独立主持人站在空白侧边；单页人物数量必须控制在 1 到 3 个之间，且不得出现 4 个及以上同角色人物铺满全页。\n"
        "23. 同一选题内，相邻页和近邻页不要复用近似动作、近似表情或近似主持姿态；正文页默认压低 large-focus，优先 small-accent 或 medium-support，继续参考 handoff 里的 role_variation_guard 和 role_variation_reason。\n"
        f"24. 最终交付图只认当前配置归档根目录下的选题目录：{handoff['archive_dir']}，不认任何临时目录或缓存图。\n"
        "25. 即使只是单独重做某页、单独重做某张封面，或单独生成某一张图片，也必须最终归档到该选题目录里的正式文件名；若先产生中间图，复制并校验正式文件存在后必须删除这张对应缓存图。\n"
        "26. 如需只重做某页，重新运行 build_ip_ppt.py 并传 --pages。\n"
        "27. 只有当本次目标页和本次目标封面都已出现在正式归档目录，且中间图片已清理，这一套图才算生成完成。\n"
        "28. 只有在用户明确要求 PPT 时，才在全部目标页都已写回 output_image 后，再传 --assemble-only 和 --out 重新装配。\n"
    )
    write_text(workflow_path, workflow)
    face_anchor_image = pathlib.Path(
        next(
            (
                item.get("role_face_anchor_image")
                for item in items
                if item.get("role_face_anchor_image")
            ),
            "",
        )
    )
    prompt_path, job_path = write_codex_thread_job_files(
        deck_title=str(next((spec.get("deck_title") for spec in specs if spec.get("deck_title")), "") or specs[0].get("page_title") or "").strip(),
        workdir=workdir,
        archive_dir=page_paths[0].parent if page_paths else workdir,
        handoff_path=handoff_path,
        specs_path=workdir / "pages-spec.json",
        workflow_path=workflow_path,
        reference_image=face_anchor_image,
        selected_pages=selected_pages,
        include_covers=is_full_run,
    )
    handoff["codex_render_thread_prompt_file"] = str(prompt_path)
    handoff["codex_render_job_file"] = str(job_path)
    handoff["codex_render_thread_required"] = True
    handoff["codex_render_current_thread_allowed"] = True
    write_json(handoff_path, handoff)
    return handoff_path, workflow_path, prompt_path, job_path


def apply_deck_title(pages: list[dict[str, Any]], deck_title: str | None) -> list[dict[str, Any]]:
    if not pages:
        return pages
    if not deck_title:
        return pages
    first = pages[0]
    if first.get("title_from_bold"):
        first["deck_title"] = deck_title
        for idx in range(1, len(pages)):
            pages[idx]["deck_title"] = deck_title
        return pages
    original_title = first["title"]
    hook = original_title if original_title != deck_title else ""
    updated = dict(first)
    updated["deck_title"] = deck_title
    updated["title"] = deck_title
    if hook:
        updated["page_hook"] = hook
    pages[0] = updated
    for idx in range(1, len(pages)):
        pages[idx]["deck_title"] = deck_title
    return pages


def normalize_page_images(page_paths: list[pathlib.Path], width: int, height: int) -> None:
    for page_path in page_paths:
        with Image.open(page_path) as img:
            if img.size == (width, height):
                continue
            rgb = img.convert("RGB")
            resized = rgb.resize((width, height))
            resized.save(page_path)


def main() -> None:
    ap = argparse.ArgumentParser(description="把中文文案生成成带固定角色的整页图片；优先支持 --- 手动分页，无 --- 时自动分页")
    ap.add_argument("--input", required=False, help="输入文案文件；如果路径不存在，可配合 --deck-title 回填现有工作包")
    ap.add_argument("--deck-title", default=None, help="可选选题标题；当 --input 缺失或路径不存在时，按标题优先回填现有工作包")
    ap.add_argument("--out", default=None, help="可选 PPTX 路径；默认不输出 PPT，只有在 --assemble-only 时才需要")
    ap.add_argument("--character", default=None, help="显式指定角色 slug")
    ap.add_argument("--render-mode", choices=["codex"], default="codex", help="只允许 codex：生成给 Codex 使用的出图工作清单")
    ap.add_argument("--workdir", default=None, help="可选工作包目录；默认固定为 SKILL_ROOT/outputs/<选题名>-work")
    ap.add_argument("--pages", default=None, help="可选页码范围，例如 1-2,4,6-7；用于只重做部分页面")
    ap.add_argument("--assemble-only", action="store_true", help="只根据现有 pages/*.png 装配 PPT，不重新生成规格与 prompt")
    args = ap.parse_args()

    out_path = pathlib.Path(args.out).expanduser().resolve() if args.out else None
    cfg = load_config()
    character_slug = resolve_default_character(args.character)
    role = load_character(character_slug)
    archive_root = resolve_archive_root(cfg)
    input_path = pathlib.Path(args.input).expanduser().resolve() if args.input else pathlib.Path("<missing-input>")
    raw_text = ""
    content_text = ""
    deck_title = (args.deck_title or "").strip() or None
    input_exists = bool(args.input) and input_path.exists()

    if input_exists:
        raw_text = input_path.read_text(encoding="utf-8")
        parsed_title, content_text = extract_deck_title(raw_text)
        deck_title = parsed_title or deck_title
    elif not deck_title:
        raise FileNotFoundError(
            f"输入文案不存在: {input_path}。若要直接复用现有工作包，请补充 --deck-title。"
        )

    archive_dir = resolve_archive_dir(deck_title, input_path, archive_root)
    outputs_root = SKILL_ROOT / "outputs"

    if args.workdir:
        workdir = pathlib.Path(args.workdir).expanduser().resolve()
    elif not input_exists:
        matched_workdir = find_existing_workdir_by_title(deck_title, outputs_root)
        if matched_workdir is None:
            raise FileNotFoundError(
                f"输入文案不存在，且未找到与选题标题匹配的现有工作包: {deck_title}"
            )
        validate_existing_workdir(matched_workdir)
        handoff_path, workflow_path, specs_path, thread_prompt_path, thread_job_path = rebind_existing_work_package_outputs(
            matched_workdir,
            archive_dir,
        )
        handoff = json.loads(handoff_path.read_text(encoding="utf-8"))
        print(f"已按选题标题回填现有工作包: {matched_workdir}")
        print(json.dumps({
            "default_character": character_slug,
            "deck_title": deck_title,
            "workdir": str(matched_workdir),
            "work_package_dir": str(matched_workdir),
            "pages_spec": str(specs_path),
            "handoff": str(handoff_path),
            "workflow": str(workflow_path),
            "codex_render_thread_prompt": str(thread_prompt_path) if thread_prompt_path else "",
            "codex_render_job": str(thread_job_path) if thread_job_path else "",
            "archive_root": str(archive_root),
            "archive_dir": str(archive_dir),
            "final_archive_dir": str(archive_dir),
            "mode": "reuse-existing-work-package",
            "reused_existing_work_package": True,
            "selected_pages": handoff.get("selected_pages", []),
            "covers": [cover.get("cover_type") for cover in handoff.get("covers", [])],
            "output": "images-only",
        }, ensure_ascii=False, indent=2))
        return
    else:
        title_folder = resolve_title_folder_name(deck_title, input_path)
        workdir = (outputs_root / f"{title_folder}-work").resolve()
    prompts_dir = ensure_dir(workdir / "prompts")
    pages_dir = ensure_dir(archive_dir)
    specs_path = workdir / "pages-spec.json"

    if args.assemble_only:
        if out_path is None:
            raise ValueError("--assemble-only 时必须传 --out；默认最终交付只输出页面图，不自动输出 PPT")
        specs = load_specs_from_workdir(workdir)
        page_paths = [pages_dir / f"page-{int(spec['page_index']):02d}.png" for spec in specs]
        ensure_archive_writeback_complete(page_paths, pages_dir)
        normalize_page_images(
            page_paths,
            int(cfg.get("default_slide_width_px", 1600)),
            int(cfg.get("default_slide_height_px", 900)),
        )
        build_ppt_from_pages(page_paths, out_path)
        print(f"已根据现有页面图片装配 PPT: {out_path}")
        print(json.dumps({
            "default_character": character_slug,
            "pages": len(specs),
            "workdir": str(workdir),
            "pages_spec": str(specs_path),
            "archive_root": str(archive_root),
            "archive_dir": str(archive_dir),
            "mode": "assemble-only",
        }, ensure_ascii=False, indent=2))
        return

    pages = parse_pages(content_text)
    pages = apply_deck_title(pages, deck_title)
    specs = [build_page_spec(page, len(pages), pages) for page in pages]
    specs = apply_role_variation_across_specs(specs)
    cover_outputs = build_cover_outputs(specs, deck_title, archive_dir, input_path)
    attach_cover_role_chain(cover_outputs, specs, role)
    write_json(specs_path, {
        "pages": specs,
        "cover_outputs": cover_outputs,
    })
    selected_pages = parse_page_selection(args.pages, len(specs))
    is_full_run = len(selected_pages) == len(specs)

    prompt_paths: list[pathlib.Path] = []
    page_paths: list[pathlib.Path] = []
    for spec in specs:
        idx = int(spec["page_index"])
        prompt_path = prompts_dir / f"page-{idx:02d}.md"
        page_path = pages_dir / f"page-{idx:02d}.png"
        if idx in selected_pages:
            write_text(prompt_path, build_prompt(spec, role))
        prompt_paths.append(prompt_path)
        page_paths.append(page_path)
    cover_prompt_paths: dict[str, pathlib.Path] = {}
    if is_full_run:
        for cover_spec in cover_outputs:
            cover_prompt_path = prompts_dir / f"{cover_spec['cover_type']}.md"
            write_text(
                cover_prompt_path,
                build_cover_prompt(
                    cover_spec,
                    role,
                ),
            )
            cover_prompt_paths[cover_spec["cover_type"]] = cover_prompt_path

    handoff_path, workflow_path, thread_prompt_path, thread_job_path = build_codex_handoff(
        specs=specs,
        cover_outputs=cover_outputs,
        selected_pages=selected_pages,
        workdir=workdir,
        prompt_paths=prompt_paths,
        page_paths=page_paths,
        cover_prompt_paths=cover_prompt_paths,
        role=role,
    )
    print(f"已生成 Codex 出图工作清单: {handoff_path}")
    print(f"已生成 Codex 操作说明: {workflow_path}")

    print(json.dumps({
        "default_character": character_slug,
        "pages": len(specs),
        "selected_pages": selected_pages,
        "covers": [cover["cover_type"] for cover in cover_outputs] if is_full_run else [],
        "workdir": str(workdir),
        "work_package_dir": str(workdir),
        "pages_spec": str(specs_path),
        "codex_render_thread_prompt": str(thread_prompt_path),
        "codex_render_job": str(thread_job_path),
        "archive_root": str(archive_root),
        "archive_dir": str(archive_dir),
        "final_archive_dir": str(archive_dir),
        "mode": args.render_mode,
        "output": "images-only",
    }, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
